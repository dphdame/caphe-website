#!/usr/bin/env python3
"""
Build: "The Value of Public Health: Measuring What Prevention Is Worth"
CAPHE practitioner-education webinar deck (comprehensive, visual, with notes).

Visuals reuse the CAPHE CEA/ROI webinar concept illustrations
(outputs/presentations/images/) and the real working-paper figures
(staged in /tmp/phv_figs). Speaker notes are first-person-plural talking
points (style-pass --social compliant).

Style: CAPHE brand template (off-white bg + navy/orange accent bars, Calibri).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image

FONT = "Calibri"
NAVY = RGBColor(0x00, 0x30, 0x80)
ORANGE = RGBColor(0xDA, 0x77, 0x0D)
OFFWHITE = RGBColor(0xFB, 0xF5, 0xEA)  # warm cream
DARK = RGBColor(0x1C, 0x1C, 0x1C)
GRAY = RGBColor(0x42, 0x42, 0x42)

IMGDIR = "/Users/victoriaperez/Projects/CAPHE/outputs/presentations/images"
FIGDIR = "/tmp/phv_figs"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _rect(slide, l, t, w, h, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def brand(slide, title_slide=False):
    _rect(slide, 0, 0, 13.333, 7.5, OFFWHITE)
    if title_slide:
        _rect(slide, 0, 0, 13.333, 0.15, NAVY)
        _rect(slide, 0, 7.35, 13.333, 0.15, ORANGE)
    else:
        _rect(slide, 0, 0, 13.333, 0.12, ORANGE)


def _set(p, size, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.name = FONT
    p.font.color.rgb = color
    p.alignment = align


def box(slide, l, t, w, h, anchor=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    return tf


def add_img(slide, name, l, t, w):
    path = name if os.path.isabs(name) else os.path.join(IMGDIR, name)
    slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))


def notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text.strip()


def title_block(slide, text, width=11.9):
    tf = box(slide, 0.7, 0.45, width, 1.05)
    _set(tf.paragraphs[0], 28, bold=True, color=NAVY)
    tf.paragraphs[0].text = text


def footnote(slide, text):
    tf = box(slide, 0.7, 6.95, 11.9, 0.45)
    _set(tf.paragraphs[0], 13, color=GRAY)
    tf.paragraphs[0].text = text


def bullets(slide, items, top=1.65, width=11.6, size=20):
    tf = box(slide, 0.9, top, width, 5.1)
    for i, (text, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if lvl == 0:
            p.text = "•  " + text
            _set(p, size, color=DARK)
            p.space_after = Pt(10)
        else:
            p.text = "       –  " + text
            _set(p, size - 1, color=GRAY)
            p.space_after = Pt(6)


# ---- slide constructors ----------------------------------------------------

def slide_title(title, subtitle, speaker, date, image=None, nt=None):
    s = prs.slides.add_slide(BLANK)
    brand(s, title_slide=True)
    tf = box(s, 1.0, 0.85, 11.33, 0.8)
    _set(tf.paragraphs[0], 22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tf.paragraphs[0].text = "California Association of Public Health Economists"
    tf = box(s, 1.0, 1.7, 11.33, 1.7)
    _set(tf.paragraphs[0], 38, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tf.paragraphs[0].text = title
    p = tf.add_paragraph(); _set(p, 24, color=GRAY, align=PP_ALIGN.CENTER)
    p.text = subtitle
    if image:
        add_img(s, image, 5.47, 3.35, 2.4)
    tf = box(s, 1.0, 5.95, 11.33, 0.9)
    _set(tf.paragraphs[0], 20, color=GRAY, align=PP_ALIGN.CENTER)
    tf.paragraphs[0].text = speaker
    p = tf.add_paragraph(); _set(p, 18, color=GRAY, align=PP_ALIGN.CENTER)
    p.text = date
    notes(s, nt)


def slide_section(part, title, nt=None):
    s = prs.slides.add_slide(BLANK)
    brand(s)
    tf = box(s, 1.0, 2.6, 11.33, 0.8)
    _set(tf.paragraphs[0], 22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tf.paragraphs[0].text = part
    tf = box(s, 1.0, 3.5, 11.33, 1.4)
    _set(tf.paragraphs[0], 40, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    tf.paragraphs[0].text = title
    notes(s, nt)


def slide_content(title, items, foot=None, size=20, image=None, nt=None):
    s = prs.slides.add_slide(BLANK)
    brand(s)
    if image:
        title_block(s, title, width=11.9)
        bullets(s, items, width=6.9, size=size)
        add_img(s, image, 8.05, 1.75, 4.6)
    else:
        title_block(s, title)
        bullets(s, items, size=size)
    if foot:
        footnote(s, foot)
    notes(s, nt)


def slide_keyconcept(label, big, example, foot=None, image=None, nt=None):
    s = prs.slides.add_slide(BLANK)
    brand(s)
    if image:
        kl, kw, al = 0.8, 6.9, PP_ALIGN.LEFT
    else:
        kl, kw, al = 1.0, 11.33, PP_ALIGN.CENTER
    tf = box(s, kl, 1.0, kw, 0.6)
    _set(tf.paragraphs[0], 20, bold=True, color=NAVY, align=al)
    tf.paragraphs[0].text = label
    tf = box(s, kl if image else 1.2, 1.9, kw if image else 10.93, 2.4,
             anchor=MSO_ANCHOR.MIDDLE)
    _set(tf.paragraphs[0], 28 if image else 30, bold=True, color=DARK, align=al)
    tf.paragraphs[0].text = big
    tf = box(s, kl if image else 1.5, 4.5, kw if image else 10.33, 2.1)
    for i, line in enumerate(example):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set(p, 19 if image else 20, color=GRAY, align=al)
        p.text = line
        p.space_after = Pt(8)
    if image:
        add_img(s, image, 8.05, 1.9, 4.6)
    if foot:
        footnote(s, foot)
    notes(s, nt)


def slide_metrics(title, cards, foot=None, image=None, nt=None):
    s = prs.slides.add_slide(BLANK)
    brand(s)
    title_block(s, title)
    span = 11.6 if not image else 6.9
    n = len(cards)
    gap = 0.4
    cw = (span - gap * (n - 1)) / n
    left0 = 0.9
    for i, (num, lab) in enumerate(cards):
        l = left0 + i * (cw + gap)
        tf = box(s, l, 2.2, cw, 1.3, anchor=MSO_ANCHOR.MIDDLE)
        _set(tf.paragraphs[0], 38, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        tf.paragraphs[0].text = num
        tf = box(s, l, 3.6, cw, 2.7)
        for j, line in enumerate(lab):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            _set(p, 18, color=DARK, align=PP_ALIGN.CENTER)
            p.text = line
            p.space_after = Pt(6)
    if image:
        add_img(s, image, 8.05, 1.95, 4.6)
    if foot:
        footnote(s, foot)
    notes(s, nt)


def slide_worked(title, setup, calc_lines, result, interp, foot=None, nt=None):
    s = prs.slides.add_slide(BLANK)
    brand(s)
    title_block(s, title)
    tf = box(s, 0.9, 1.55, 11.6, 0.9)
    _set(tf.paragraphs[0], 19, color=GRAY)
    tf.paragraphs[0].text = setup
    tf = box(s, 0.9, 2.5, 7.2, 3.0)
    for i, line in enumerate(calc_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set(p, 20, color=DARK)
        p.text = line
        p.space_after = Pt(8)
    tf = box(s, 8.4, 2.5, 4.1, 1.4, anchor=MSO_ANCHOR.MIDDLE)
    _set(tf.paragraphs[0], 26, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    tf.paragraphs[0].text = result
    tf = box(s, 8.4, 4.1, 4.1, 2.4)
    _set(tf.paragraphs[0], 17, color=GRAY)
    tf.paragraphs[0].text = interp
    if foot:
        footnote(s, foot)
    notes(s, nt)


def slide_figure(title, fig, caption=None, nt=None, width=8.0):
    s = prs.slides.add_slide(BLANK)
    brand(s)
    title_block(s, title)
    path = os.path.join(FIGDIR, fig)
    w, h = Image.open(path).size
    iw = width
    ih = iw * h / w
    if ih > 5.35:
        ih = 5.35
        iw = ih * w / h
    left = (13.333 - iw) / 2
    s.shapes.add_picture(path, Inches(left), Inches(1.55), width=Inches(iw))
    if caption:
        tf = box(s, 0.9, 1.55 + ih + 0.05, 11.5, 0.5)
        _set(tf.paragraphs[0], 15, color=GRAY, align=PP_ALIGN.CENTER)
        tf.paragraphs[0].text = caption
    notes(s, nt)


# ============================================================================
# DECK
# ============================================================================

slide_title(
    "The Value of Public Health",
    "Measuring What Prevention Is Worth",
    "Victoria Cholette, PhD",
    "Free Webinar Series  ·  caphegroup.org",
    image="cea-hero-california.png",
    nt="Welcome. Today we ask what public health is worth. We also ask how we "
       "put numbers on prevention. We share new California evidence from our "
       "own work. Then we walk through everyday tools for weighing a program. "
       "We hope to leave knowing which method answers which question.",
)

slide_content("Today's Objectives", [
    ("Share new California evidence on prevention's value", 0),
    ("our replication and extension of Brown (2016)", 1),
    ("Understand why measuring that value is hard", 0),
    ("Tour the toolkit economists use:", 0),
    ("common health units, cost-effectiveness, ROI", 1),
    ("budget impact, comparators, sensitivity analysis", 1),
    ("See worked examples from California public health", 0),
    ("Know how to choose a method for our own work", 0),
], nt="Here is our path for the hour. We start with the evidence, including a "
      "study we ran on California counties. Then we ask why prevention is so "
      "hard to value. Most of our time goes to the toolkit, with worked "
      "examples we can borrow. We close on how to match a method to a question.")

# ---- Part 1 ----------------------------------------------------------------

slide_keyconcept(
    "THE CORE QUESTION",
    "Prevention's payoff is the bad outcome that never happens.",
    ["How do we value an illness that never occurred,",
     "a hospitalization avoided, a life-year preserved?",
     "This is the measurement problem at the heart of our field."],
    image="cea-concept-question.png",
    nt="Prevention works by making bad things not happen. That is what makes "
       "it satisfying and also what makes it hard to measure. We cannot count "
       "the heart attack that never came. So we lean on comparisons and models "
       "to estimate the outcome we avoided.")

slide_metrics("The Headline Evidence: Brown (2016)", [
    ("$67–88", ["Return per dollar of", "public health spending"]),
    ("2001–2009", ["Study period,", "40 CA counties"]),
    ("IV design", ["Corrects reverse causation", "& omitted variables", "(Koyck lag model)"]),
], foot="Brown TT. Returns on Investment in CA County Departments of "
        "Public Health. Am J Public Health 2016;106(8):1477–1482.",
    image="roi-concept-multiplier.png",
    nt="The most cited number in our space comes from Brown. He estimated a "
       "return of sixty-seven to eighty-eight dollars per dollar of county "
       "public health spending. The figure is large, and it rests on a careful "
       "instrumental-variables design rather than a simple comparison. We "
       "build on this work later.")


slide_content("Why Measuring Prevention's Value Is Hard", [
    ("The counterfactual is invisible", 0),
    ("we count events that did not happen", 1),
    ("Benefits are diffuse and delayed", 0),
    ("costs come now; payoffs arrive years later", 1),
    ("\"Public health\" spending is mixed", 0),
    ("clinical care, admin, emergency response", 1),
    ("prevention is 40–60% of county budgets", 1),
    ("Lumping it together can understate ROI", 0),
], image="roi-concept-investment.png",
    nt="Four things make this hard. We cannot observe the outcome we prevented. "
       "The benefits show up slowly while the bill arrives now. The label "
       "public health covers many activities, so the spending is mixed. And "
       "when we lump prevention in with everything else, we may understate "
       "what prevention alone returns.")

# ---- Part 2: Our Research --------------------------------------------------
slide_section("Our Research", "New California Evidence",
    nt="Now to our own work. We set out to test whether Brown's result "
       "replicates. We also asked what it misses. The answers surprised us in "
       "a few places.")

slide_content("What We Did: Replicate and Extend Brown", [
    ("California, all 57 counties, 2003–2023", 0),
    ("1,197 county-year observations", 1),
    ("Same method as Brown: Lewbel (2012) IV", 0),
    ("internal instruments from heteroskedasticity", 1),
    ("no hard-to-find external instrument needed", 1),
    ("Outcome: all-cause mortality per 100,000", 0),
    ("Built from detailed county financial reports", 0),
    ("validated with county health officers", 1),
], foot="Our working paper (2025); replicates and extends Brown (2016).",
    image="cea-hero-california.png",
    nt="We assembled a panel of all fifty-seven California counties over two "
       "decades, about twelve hundred county-years. We used the same Lewbel "
       "method as Brown, which builds its instruments from the data itself. "
       "We measured spending from detailed county financial reports and "
       "checked our classification with county health officers.")

slide_metrics("We Reproduced Brown — and the Bias He Corrected", [
    ("−9.16", ["Our Lewbel IV estimate", "Brown (2016): −9.1", "a close match"]),
    ("0.18", ["Plain OLS estimate", "null (p = 0.64)", "wrong sign"]),
    ("50×", ["Gap between OLS", "and IV — the size", "of the bias"]),
], foot="Reverse causation: mortality shocks raise spending later, hiding the benefit in OLS.",
    nt="Our central estimate landed at minus nine point one six. That is almost "
       "exactly Brown's minus nine point one. Plain regression looks far "
       "different. Ordinary least squares gives a null with the wrong sign, a "
       "gap of about fifty times. That gap is the reverse-causation bias, since "
       "mortality shocks tend to raise spending later.")

slide_figure("The Bias, Made Visible", "figure_3_ols_vs_iv.png",
    caption="Plain OLS shows no effect; Lewbel IV reveals a large mortality reduction.",
    nt="Here is the same point in one picture. On the left, ordinary "
       "regression sits at zero. On the right, the method that handles "
       "endogeneity shows a clear mortality reduction with tight confidence "
       "bounds. Reading only the left bar, we would conclude prevention does "
       "nothing, which is the trap.")

slide_content("The Result Holds Up", [
    ("47 alternative specifications tested", 0),
    ("46 of 47 stay significant and beneficial", 1),
    ("Drop any one county: the estimate barely moves", 0),
    ("leave-one-out SD = 0.052", 1),
    ("Benefit concentrates in small counties", 0),
    ("−8.11 in small counties; not detectable in large", 1),
    ("a power limit, not an absent benefit", 1),
], nt="We pushed on the result from many angles. Across forty-seven "
      "specifications, forty-six stayed beneficial and significant. Dropping "
      "any single county barely moved the estimate. The benefit concentrates "
      "in small counties. In large ones we cannot detect it, which looks like "
      "a sample-size limit, not a true absence.")


slide_keyconcept(
    "THE MAIN FINDING",
    "Spending prevented about twice as much mortality in crisis years.",
    ["Effect excluding crisis years ≈4.6; including them ≈9.2 —",
     "about double. Standard analyses measure only normal",
     "times, so they may understate the value by 50–100%."],
    foot="The option value of capacity: it pays off most when a crisis hits.",
    image="roi-concept-prevention-pays.png",
    nt="Our headline finding is about crisis value. The mortality effect of "
       "spending roughly doubled once we include the pandemic years. Standard "
       "analyses look only at normal times, so they may understate the value "
       "of capacity by half to fully double. We think of this as the option "
       "value of having public health infrastructure ready.")

slide_figure("The Option Value of Capacity", "figure_1_covid_option_value.png",
    caption="The mortality effect roughly doubles once crisis years are included.",
    nt="The chart shows the effect excluding crisis years next to the effect "
       "that includes them. The jump is close to double. Infrastructure that "
       "looks modestly useful in calm years can matter a great deal when a "
       "crisis arrives. That extra value is easy to miss.")

slide_keyconcept(
    "A SECOND FINDING",
    "Public health and healthcare access work better together.",
    ["The mortality effect was about 3× larger after the ACA",
     "(−3.07 pre-ACA  →  −8.93 post-ACA).",
     "Coverage and prevention look like complements."],
    image="slide-concept-epi-econ-complement.png",
    nt="We also found a complement. After the Affordable Care Act, the "
       "mortality effect of public health spending was about three times "
       "larger. Coverage and prevention seem to reinforce each other, so "
       "investing in both together may return more than either alone. That "
       "has direct relevance for states weighing Medicaid expansion.")

slide_figure("Stronger After the ACA", "figure_2_aca_complementarity.png",
    caption="The effect is about 3× larger post-ACA, suggesting complementarity.",
    nt="Here the pre-ACA effect sits next to the post-ACA effect, and the "
       "difference is roughly threefold. We read this as evidence that "
       "prevention needs a working path to care to translate into saved "
       "lives. The two investments appear to multiply rather than substitute.")

slide_metrics("The Bottom Line", [
    ("$109K", ["Public health spending", "per life saved"]),
    ("125 : 1", ["Benefit-cost ratio", "at a $13.6M", "value of a life"]),
    ("50–100%", ["How much standard", "analyses may", "understate the value"]),
], foot="Value-of-statistical-life standard, HHS 2025. Crisis value raises returns further.",
    image="roi-concept-returns.png",
    nt="About one dollar per capita in added spending is linked to roughly "
       "nine fewer deaths per 100,000. Put in dollar terms, that is about one "
       "hundred nine thousand dollars per life saved. Against the federal value "
       "of a statistical life, the benefit-cost ratio is near one hundred "
       "twenty-five to one. And because standard analyses miss crisis value, "
       "they may understate the return by half to fully double.")

# ---- Part 3 ----------------------------------------------------------------
slide_section("The Toolkit", "Measuring Prevention's Value",
    nt="Numbers like these come from a way of thinking. Before any formula, "
       "economists start with one idea: every choice has a cost we cannot see.")

slide_content("\"It Works\" Isn't Enough", [
    ("Showing a program improves outcomes is step one", 0),
    ("The real question is what we give up to fund it", 0),
    ("opportunity cost: the next best use of the dollar", 1),
    ("Every yes is a no to something else", 0),
    ("A program can work and still be the wrong choice", 0),
    ("if another option buys more health per dollar", 1),
], foot="Methods Lab: Why \"It Works\" Isn't Enough",
    image="slide-concept-opportunity-cost.png",
    nt="Proving a program works is only the first step. The harder question is "
       "what we gave up to fund it, the opportunity cost. Every dollar we "
       "commit is a dollar unavailable elsewhere. So a program can work and "
       "still be the wrong call if another option would buy more health.")


# ---- Part 4 ----------------------------------------------------------------

slide_content("The Comparison Problem", [
    ("Program A prevents 100 deaths", 0),
    ("Program B prevents 500 hospitalizations", 0),
    ("Program C improves quality of life for 1,000", 0),
    ("Which delivers the most health?", 0),
    ("We need a common unit capturing both", 0),
    ("mortality (length of life) and morbidity (quality)", 1),
], image="cea-concept-measures.png",
    nt="Take three programs. One prevents deaths, one prevents "
       "hospitalizations, one improves daily quality of life. Which gives the "
       "most health? We cannot say until we put them in one unit. That unit "
       "has to capture both how long we live and how well we live.")

slide_keyconcept(
    "KEY CONCEPT",
    "One QALY = one year of life in perfect health.",
    ["A year at reduced quality counts as less than 1 QALY.",
     "Extend life 5 years at 0.6 quality  →  5 × 0.6 = 3 QALYs.",
     "QALYs capture quantity and quality of life together."],
    foot="DALYs are a related measure: years of healthy life lost.",
    nt="The quality-adjusted life year is our common unit. One QALY is a year "
       "in full health, and a year at lower quality counts as less. So five "
       "years at sixty percent quality equals three QALYs. The measure lets us "
       "fold length and quality of life into a single number.")

slide_content("QALYs: Useful, and Contested", [
    ("Standard in research; we will see them cited", 0),
    ("Criticism: may undervalue care for the disabled", 0),
    ("or for elderly populations", 1),
    ("a life-year with disability counted as less", 1),
    ("Some read this as bias against disability", 0),
    ("Alternatives weight all life-years equally", 0),
    ("Use thoughtfully; know the debate", 0),
], foot="An active area of ethical and methodological discussion.",
    nt="QALYs are standard, and they are contested. The main worry is that "
       "they can count a year lived with disability as worth less, which some "
       "read as bias. Alternatives weight every life-year equally. We will not "
       "settle that debate today. The aim is to use the measure while knowing "
       "the critique.")


# ---- Part 5 ----------------------------------------------------------------

slide_content("What CEA Is (and Isn't)", [
    ("IS a tool to compare value across programs", 0),
    ("IS a way to summarize cost and outcome in one ratio", 0),
    ("IS NOT a decision-maker; it informs, doesn't decide", 0),
    ("IS NOT a judgment of worth; it measures efficiency", 0),
    ("IS NOT the only input (equity, feasibility, politics)", 0),
], image="slide-concept-cost-effectiveness.png",
    nt="Let us be clear about what the tool is. Cost-effectiveness analysis "
       "compares value across programs and sums it into one ratio. It does not "
       "make the decision for us, and it does not judge human worth. It "
       "measures efficiency, which is one input among equity, feasibility, "
       "and politics.")

slide_keyconcept(
    "THE CORE METRIC",
    "ICER = (Cost_new − Cost_now) ÷ (Effect_new − Effect_now)",
    ["The cost of buying one more unit of health.",
     "\"How much extra do we pay for each extra unit gained?\"",
     "The comparison is always against current practice."],
    image="cea-concept-icer-ratio.png",
    nt="The core metric is the incremental cost-effectiveness ratio. It is the "
       "extra cost divided by the extra health, against current practice. In "
       "plain terms, how much more do we pay for each additional unit of "
       "health. That ratio is what lets us rank options on the same scale.")

slide_worked(
    "Worked Example: Diabetes Screening",
    "A county weighs expanding diabetes screening to a new population.",
    ["Current:  $500K  →  50 cases prevented",
     "Expanded: $800K  →  80 cases prevented",
     "──────────────────────────",
     "Extra cost:   $300,000",
     "Extra effect: 30 cases",
     "ICER = $300K ÷ 30"],
    "$10,000\nper case prevented",
    "We pay $10,000 for each additional case prevented. Worth it? It depends "
    "on the downstream cost of a diabetes case and the threshold we apply.",
    nt="Here is the ratio in action. Expanding screening costs three hundred "
       "thousand more and prevents thirty more cases, so the incremental cost "
       "is ten thousand dollars per case. Whether that is worth it depends on "
       "what a diabetes case costs us downstream and on the threshold we hold.")

slide_worked(
    "California Example: Promotora Case Management",
    "Diabetes case management with nurse managers and promotora peer educators "
    "at community health centers (low-income, predominantly Latino patients).",
    ["Cost: $1,537/patient, year 1",
     "Medi-Cal cohort (n=1,213):",
     "  HbA1c:  −0.5 % points",
     "  Systolic BP:  −1.9 mmHg",
     "──────────────────────────",
     "ICER (Medi-Cal cohort)"],
    "$44,941\nper QALY gained",
    "Just below a $50,000/QALY benchmark. The uninsured cohort reached "
    "$10,141/QALY. Cohort differences can reflect baseline severity and "
    "sample composition, not the program alone.",
    nt="A California example grounds it. A promotora-led diabetes program cost "
       "about forty-five thousand dollars per QALY in the Medi-Cal group, just "
       "under a common benchmark. The uninsured group came in far lower. We "
       "read that gap with care, since cohorts differ in baseline severity, "
       "not only in program effect.")

slide_metrics("Decision Thresholds: What's \"Cost-Effective\"?", [
    ("< $50K", ["per QALY", "Highly", "cost-effective", "Strong case"]),
    ("$50–150K", ["per QALY", "Potentially", "cost-effective", "Context-dependent"]),
    ("> $150K", ["per QALY", "Less", "cost-effective", "Higher bar"]),
], foot="Guidelines, not rules. The U.S. has no official threshold; some argue $200K+.",
    image="cea-concept-thresholds.png",
    nt="So what counts as cost-effective? A rough convention puts under fifty "
       "thousand per QALY as strong value. Fifty to one hundred fifty thousand "
       "is context-dependent. Above that sets a higher bar. These are "
       "guidelines, not rules. The United States has no official threshold, "
       "and some argue it should be higher.")

slide_content("When \"Not Cost-Effective\" Isn't the End", [
    ("Cost-effectiveness is one input, not the verdict", 0),
    ("Equity: serving the underserved can justify more", 0),
    ("Legal or ethical duty: some services are required", 0),
    ("Community priorities: values matter beyond efficiency", 0),
    ("No alternative: an expensive option beats nothing", 0),
    ("Strategic: capacity now pays off in ways CEA misses", 0),
], image="cea-concept-beyond.png",
    nt="A high ratio does not end the conversation. Equity can justify "
       "spending more on the underserved. Some services are legally or "
       "ethically required. Community values matter. Sometimes no alternative "
       "exists, and an expensive option still beats nothing. And capacity can "
       "pay off in ways the ratio never captures.")

# ---- Part 6 ----------------------------------------------------------------


slide_content("Budget Impact: Can We Afford It?", [
    ("Cost-effective is not the same as affordable", 0),
    ("Budget impact asks a different question:", 0),
    ("what does this cost our budget, total, per year?", 1),
    ("A program can be great value yet unaffordable", 0),
    ("or affordable yet poor value", 1),
    ("Economists ask both questions, separately", 0),
], foot="Methods Lab: Budget Impact: Can We Afford It?",
    image="roi-concept-investment.png",
    nt="Cost-effective and affordable are not the same thing. Budget impact "
       "asks what a program costs our budget in total, year by year. A program "
       "can be excellent value and still break the budget, or be cheap and "
       "poor value. We ask both questions, and we keep them separate.")

slide_content("ROI and Benefit-Cost Analysis", [
    ("CEA keeps health in health units (QALYs, cases)", 0),
    ("ROI / benefit-cost puts benefits in dollars", 0),
    ("dollars returned per dollar spent (e.g., $67–88)", 1),
    ("Strength: one currency, easy to communicate", 0),
    ("Cost: monetizing health requires assumptions", 0),
    ("how much is a life-year worth in dollars?", 1),
    ("Choose the framing the audience trusts", 0),
], image="roi-concept-audiences.png",
    nt="Cost-effectiveness keeps health in health units. Return on investment "
       "and benefit-cost analysis convert benefits into dollars, like the "
       "sixty-seven to eighty-eight dollar figure. Dollars are easy to "
       "communicate, but turning health into money takes assumptions. We pick "
       "the framing our audience will trust.")

# ---- Part 7 ----------------------------------------------------------------

slide_content("The Foundation: Does It Actually Cause the Outcome?", [
    ("Every ROI and ICER rests on a causal claim", 0),
    ("the program caused the improvement", 1),
    ("Study designs differ in how well they isolate that", 0),
    ("RCTs and strong quasi-experiments at the top", 1),
    ("simple before-after near the bottom", 1),
    ("Common traps: confounding, selection,", 0),
    ("regression to the mean, the before-after trap", 1),
], foot="Methods Lab: Study Design Hierarchy for Causal Inference",
    nt="Under every ratio is a causal claim: that the program caused the "
       "change. Study designs differ in how well they support that claim. "
       "Randomized trials and strong quasi-experiments sit near the top. "
       "Simple before-after sits near the bottom. Confounding, selection, and "
       "regression to the mean are the usual traps.")



# ---- In Practice -----------------------------------------------------------




# ---- Close -----------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
brand(s)
add_img(s, "cea-cta-community.png", 5.27, 1.0, 2.8)
tf = box(s, 1.0, 4.0, 11.33, 1.2, anchor=MSO_ANCHOR.MIDDLE)
_set(tf.paragraphs[0], 40, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tf.paragraphs[0].text = "Thank You"
tf = box(s, 1.0, 5.2, 11.33, 1.6)
for i, line in enumerate([
    "Questions?",
    "caphegroup.org  ·  info@caphegroup.org",
    "Methods Lab  ·  Monthly peer review sessions",
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    _set(p, 22 if i == 0 else 20, color=GRAY, align=PP_ALIGN.CENTER)
    p.text = line
    p.space_after = Pt(8)
notes(s, "Thanks for joining. The slides, the Methods Lab, and our working "
         "paper are all linked at caphegroup.org. We welcome questions now. We "
         "also hope to see folks at the monthly peer review sessions, where we "
         "work through real analyses together.")

out = "/tmp/CAPHE_ValueOfPublicHealth_2026-06-18.pptx"
prs.save(out)
print("Saved:", out, "—", len(prs.slides._sldIdLst), "slides")
