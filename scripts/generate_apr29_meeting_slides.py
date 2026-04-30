#!/usr/bin/env python3
"""
CAPHE Member Meeting, April 29, 2026.
Discussion-style deck. One question per slide, room for conversation.

Agenda:
  1. Welcome (3 min)
  2. Health Economics within CDPH, Gilda (10 min)
  3. What's on your mind? round-robin (5 min)
  4. April 9 didn't land. June 24 needs your help (15 min)
  5. Summer slot, open call to members (10 min)
  6. Fall, Gilda's topic, member input (8 min)
  7. Office hours, design with us (5 min)
  8. Wrap (4 min)
"""

from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# CAPHE Brand Colors (match caphegroup.org)
PRIMARY = RGBColor(0, 48, 128)         # #003080 Institutional Blue
ACCENT = RGBColor(218, 119, 13)        # #DA770D Poppy Gold
BG_WARM = RGBColor(250, 250, 248)      # #FAFAF8
TEXT_PRIMARY = RGBColor(28, 28, 28)
TEXT_SECONDARY = RGBColor(66, 66, 66)
TEXT_MUTED = RGBColor(100, 100, 100)
WHITE = RGBColor(255, 255, 255)

SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent
OUTPUT_DIR = PROJECT_DIR / "outputs" / "presentations"
MEETING_DATE = datetime(2026, 4, 29)


def create_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_top_bar(slide, prs, color=None):
    if color is None:
        color = PRIMARY
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_bottom_bar(slide, prs, color=None):
    if color is None:
        color = ACCENT
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_section_header(slide, prs, label_text, title_text):
    """Header band with section label + title."""
    header_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.5))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = PRIMARY
    header_bar.line.fill.background()

    label = slide.shapes.add_textbox(Inches(0.75), Inches(0.3), Inches(11.833), Inches(0.4))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = label_text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    title = slide.shapes.add_textbox(Inches(0.75), Inches(0.7), Inches(11.833), Inches(0.85))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_bg(slide, prs, color=None):
    if color is None:
        color = BG_WARM
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()


# Slide 1: Title
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_top_bar(slide, prs)

    title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "CAPHE"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "California Association of Public Health Economists"
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    info = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1.5))
    tf = info.text_frame
    p = tf.paragraphs[0]
    p.text = "Member Meeting"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "April 29, 2026"
    p.font.size = Pt(26)
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    add_bottom_bar(slide, prs)


# Slide 2: Today is mostly us asking you
def slide_framing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_top_bar(slide, prs)

    title = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Today is mostly us asking you for input"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    sub = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), Inches(11.833), Inches(0.7))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Three things, plus a few quick updates."
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_SECONDARY

    items = slide.shapes.add_textbox(Inches(1.5), Inches(2.7), Inches(10.5), Inches(4))
    tf = items.text_frame
    tf.word_wrap = True

    questions = [
        ("1.", "What should CAPHE be asking CDPH for?"),
        ("2.", "What might make our webinars more useful, or more findable?"),
        ("3.", "Who wants to present this summer?"),
    ]
    for i, (num, q) in enumerate(questions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run1 = p.add_run() if i > 0 else p.runs[0] if p.runs else None
        p.text = f"{num}   {q}"
        p.font.size = Pt(28)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(20)


# Slide 3: Agenda
def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_top_bar(slide, prs)

    header = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Today's Agenda"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    content = slide.shapes.add_textbox(Inches(1), Inches(1.6), Inches(11.333), Inches(5.5))
    tf = content.text_frame
    tf.word_wrap = True

    items = [
        ("1.", "Welcome", "3 min"),
        ("2.", "Health Economics within CDPH (Gilda)", "10 min"),
        ("3.", "What's on your mind? Round-robin.", "5 min"),
        ("4.", "April 9 didn't land. June 24 needs your help.", "15 min"),
        ("5.", "Summer slot: open call to members", "10 min"),
        ("6.", "Fall: Gilda's topic, what would members find useful?", "8 min"),
        ("7.", "Office hours: design with us", "5 min"),
        ("8.", "Wrap and next steps", "4 min"),
    ]
    for i, (num, text, time) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{num}   {text}    ({time})"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(10)


# Slide 4: Item 2, Health Economics within CDPH
def slide_cdph(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_section_header(slide, prs, "ITEM 2 (GILDA)", "Health Economics within CDPH")

    sub = slide.shapes.add_textbox(Inches(0.75), Inches(1.85), Inches(11.833), Inches(0.6))
    tf = sub.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The health economics function is moving away from CDPH. Our letter was written before that shift."
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY

    # Left column: what the letter asked for
    left_label = slide.shapes.add_textbox(Inches(0.75), Inches(2.7), Inches(6), Inches(0.4))
    tf = left_label.text_frame
    p = tf.paragraphs[0]
    p.text = "WHAT THE LETTER ASKED"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    left = slide.shapes.add_textbox(Inches(0.75), Inches(3.15), Inches(6), Inches(3.5))
    tf = left.text_frame
    tf.word_wrap = True
    asks = [
        "Position CAPHE as a professional resource to CDPH.",
        "Adopt a California reference case for public health economic evaluation.",
        "Build HRQoL and cost data into modernization (BRFSS, CHIS).",
        "Engage health economists at scoping, not retrofit.",
    ]
    for i, item in enumerate(asks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(8)

    # Right column: discussion
    right_label = slide.shapes.add_textbox(Inches(7.1), Inches(2.7), Inches(5.8), Inches(0.4))
    tf = right_label.text_frame
    p = tf.paragraphs[0]
    p.text = "DOES THIS CHANGE NOW?"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    right = slide.shapes.add_textbox(Inches(7.1), Inches(3.15), Inches(5.8), Inches(3.5))
    tf = right.text_frame
    tf.word_wrap = True
    prompts = [
        "Does the letter's frame still hold?",
        "If health economics is leaving CDPH, where should CAPHE push instead?",
        "Who else at the state should know about us?",
        "What would make our next ask more concrete?",
    ]
    for i, item in enumerate(prompts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(8)


# Slide 5: Item 3, What's on your mind round-robin
def slide_whats_on_your_mind(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_section_header(slide, prs, "ITEM 3 (ROUND-ROBIN)", "What's on your mind?")

    big_q = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.833), Inches(2))
    tf = big_q.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Anything in the field that's caught your attention since February?"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    sub = slide.shapes.add_textbox(Inches(0.75), Inches(4.6), Inches(11.833), Inches(2))
    tf = sub.text_frame
    tf.word_wrap = True

    examples = [
        "CA budget / May Revise pressures on Medi-Cal?",
        "Federal Medicaid policy or data releases?",
        "A paper, dataset, or report worth flagging?",
    ]
    for i, ex in enumerate(examples):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•   " + ex
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_SECONDARY
        p.space_after = Pt(6)

    note = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(11.833), Inches(0.5))
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.text = "We don't have to solve it. Just surface it."
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED


# Slide 6: Item 4, April 9 to June 24 honest reframe
def slide_april_to_june(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_section_header(slide, prs, "ITEM 4 (JUNE 24)", "April 9 didn't land. June 24 needs your help.")

    body = slide.shapes.add_textbox(Inches(0.75), Inches(2), Inches(11.833), Inches(1.6))
    tf = body.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "April 9 had no attendees. We'd like to figure out why, together."
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_PRIMARY

    p = tf.add_paragraph()
    p.text = " "
    p.font.size = Pt(6)

    p = tf.add_paragraph()
    p.text = "June 24: Adrienne Sabety and Maya Rossin-Slater (Stanford)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    p = tf.add_paragraph()
    p.text = "Preschool entry age and developmental and behavioral health diagnoses among low-income children in Medicaid."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_SECONDARY

    q_label = slide.shapes.add_textbox(Inches(0.75), Inches(4.6), Inches(11.833), Inches(0.5))
    tf = q_label.text_frame
    p = tf.paragraphs[0]
    p.text = "WHAT MIGHT WORK BETTER?"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    qs = slide.shapes.add_textbox(Inches(0.75), Inches(5.1), Inches(11.833), Inches(2.2))
    tf = qs.text_frame
    tf.word_wrap = True
    items = [
        "Are we picking topics our people actually want?",
        "Is the format right? Length, time of day, lecture vs. discussion, recording?",
        "Who in your network would want this, and how do they usually hear about things?",
        "Partners we haven't tried: CHEAC, CCLHO, AcademyHealth CA, schools of public health?",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(6)


# Slide 7: Item 5, Summer slot
def slide_summer(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_section_header(slide, prs, "ITEM 5 (SUMMER)", "Open call: who wants to present?")

    body = slide.shapes.add_textbox(Inches(0.75), Inches(2), Inches(11.833), Inches(1))
    tf = body.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "We have an open summer slot. Member work, work-in-progress, anything you want feedback on."
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_PRIMARY

    fmt_label = slide.shapes.add_textbox(Inches(0.75), Inches(3.4), Inches(11.833), Inches(0.5))
    tf = fmt_label.text_frame
    p = tf.paragraphs[0]
    p.text = "FORMAT OPTIONS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    options = slide.shapes.add_textbox(Inches(0.75), Inches(3.9), Inches(11.833), Inches(2.2))
    tf = options.text_frame
    tf.word_wrap = True
    items = [
        "Single 30-minute talk + Q&A",
        "Two 15-minute lightning talks",
        "Works-in-progress methods clinic: bring a problem, group reacts.",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•   " + item
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(8)

    ask = slide.shapes.add_textbox(Inches(0.75), Inches(6.3), Inches(11.833), Inches(0.7))
    tf = ask.text_frame
    p = tf.paragraphs[0]
    p.text = "Hands? Nominations? Confirm by mid-May."
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = PRIMARY


# Slide 8: Item 6, Gilda fall
def slide_fall(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_section_header(slide, prs, "ITEM 6 (FALL)", "Gilda is presenting in the fall.")

    big_q = slide.shapes.add_textbox(Inches(0.75), Inches(2.3), Inches(11.833), Inches(2))
    tf = big_q.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "What would you find most useful?"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    sub = slide.shapes.add_textbox(Inches(0.75), Inches(4.4), Inches(11.833), Inches(2))
    tf = sub.text_frame
    tf.word_wrap = True
    items = [
        "What's a question you wish a state-level health economist would answer?",
        "Free public webinar or member-only?",
        "Anything Gilda should know before she finalizes the topic?",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•   " + item
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(8)


# Slide 9: Item 7, Office hours
def slide_office_hours(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_section_header(slide, prs, "ITEM 7 (OFFICE HOURS)", "Help us design the pilot.")

    body = slide.shapes.add_textbox(Inches(0.75), Inches(2), Inches(11.833), Inches(1.2))
    tf = body.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "We committed to pilot office hours in February. We have not yet because we want you to tell us the format."
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_PRIMARY

    pairs_label = slide.shapes.add_textbox(Inches(0.75), Inches(3.6), Inches(11.833), Inches(0.5))
    tf = pairs_label.text_frame
    p = tf.paragraphs[0]
    p.text = "DESIGN CHOICES"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    pairs = slide.shapes.add_textbox(Inches(0.75), Inches(4.1), Inches(11.833), Inches(2.5))
    tf = pairs.text_frame
    tf.word_wrap = True
    items = [
        "Drop-in   or   by appointment?",
        "Open Q&A   or   bring-a-problem methods clinic?",
        "Organized by topic (CEA, causal inference, data)   or   general?",
        "Audience: LHJ staff, program managers, students?",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•   " + item
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(8)


# Slide 10: Wrap
def slide_wrap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, prs)
    add_top_bar(slide, prs)

    title = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Things to take with us"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    actions = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4))
    tf = actions.text_frame
    tf.word_wrap = True
    items = [
        "Topic and format ideas for future webinars (keep them coming).",
        "Names of people or partners worth inviting for June 24.",
        "Anyone interested in presenting this summer? Confirm by mid-May.",
        "Topic suggestions for Gilda's fall talk.",
        "Office hours: format ideas to bring back next meeting.",
    ]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•   " + item
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(10)

    next_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.4), Inches(11.833), Inches(0.6))
    tf = next_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Next meeting: TBD   |   Next webinar: June 24, 2026 (Sabety & Rossin-Slater)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    add_bottom_bar(slide, prs)


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = create_presentation()

    slide_title(prs)
    slide_framing(prs)
    slide_agenda(prs)
    slide_cdph(prs)
    slide_whats_on_your_mind(prs)
    slide_april_to_june(prs)
    slide_summer(prs)
    slide_fall(prs)
    slide_office_hours(prs)
    slide_wrap(prs)

    out_path = OUTPUT_DIR / f"CAPHE_Meeting_{MEETING_DATE.strftime('%Y-%m-%d')}.pptx"
    prs.save(str(out_path))
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
