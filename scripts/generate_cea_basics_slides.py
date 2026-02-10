#!/usr/bin/env python3
"""
Generate CAPHE-branded slides for:
"Understanding Cost-Effectiveness: The Basics"
April 9, 2026 - Free Webinar

Target audience: County health department staff, policy analysts
Duration: 1 hour (12:00 PM - 1:00 PM PT)
Prerequisite: Feb 12 intro webinar concepts (or self-study)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# =============================================================================
# CAPHE Brand Colors (matching caphegroup.org)
# =============================================================================
PRIMARY = RGBColor(0x00, 0x30, 0x80)      # Institutional Blue
PRIMARY_LIGHT = RGBColor(0x00, 0x41, 0xA5)
ACCENT = RGBColor(0xDA, 0x77, 0x0D)       # Poppy Gold
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_WARM = RGBColor(0xFA, 0xFA, 0xF8)
TEXT_PRIMARY = RGBColor(0x1C, 0x1C, 0x1C)
TEXT_SECONDARY = RGBColor(0x42, 0x42, 0x42)
TEXT_MUTED = RGBColor(0x64, 0x64, 0x64)
SUCCESS = RGBColor(0x2E, 0x7D, 0x32)       # Green for "cost-effective"
WARNING = RGBColor(0xF5, 0x7C, 0x00)       # Orange for caution

# Image paths
IMAGES_DIR = Path("/Users/victoriaperez/Projects/CAPHE/07_website/outputs/presentations/images")


def create_presentation():
    """Create presentation with 16:9 widescreen dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_image_if_exists(slide, image_name, left, top, width=None, height=None):
    """Add image to slide if it exists, otherwise skip silently."""
    image_path = IMAGES_DIR / image_name
    if image_path.exists():
        if width and height:
            slide.shapes.add_picture(str(image_path), left, top, width, height)
        elif width:
            slide.shapes.add_picture(str(image_path), left, top, width=width)
        elif height:
            slide.shapes.add_picture(str(image_path), left, top, height=height)
        else:
            slide.shapes.add_picture(str(image_path), left, top)
        return True
    return False


def add_title_slide(prs, title, subtitle, date):
    """Add title slide with CAPHE branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_WARM
    bg.line.fill.background()

    # Top accent bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY
    top_bar.line.fill.background()

    # CAPHE header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CAPHE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "California Association of Public Health Economists"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.alignment = PP_ALIGN.CENTER

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.2), Inches(11.833), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.8), Inches(11.833), Inches(0.75))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.8), Inches(11.833), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = date
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER

    # Bottom accent bar
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = ACCENT
    bottom_bar.line.fill.background()


def add_section_slide(prs, title, subtitle=None):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3), Inches(11.833), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = ACCENT
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)


def add_content_slide(prs, title, bullets, note=None, image=None):
    """Add content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    content_width = Inches(7) if image else Inches(11.833)

    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), content_width, Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet.startswith("  ") or bullet.startswith("\t"):
            p.text = "    " + bullet.strip()
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.text = "• " + bullet.strip()
            p.font.size = Pt(22)
            p.level = 0

        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(12)

    if image:
        add_image_if_exists(slide, image, Inches(8.5), Inches(1.6), width=Inches(4.5))

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(11.833), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEXT_MUTED


def add_key_concept_slide(prs, concept, definition, example=None, image=None):
    """Add a slide highlighting a key concept."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    text_width = Inches(7) if image else Inches(11.833)

    label_box = slide.shapes.add_textbox(Inches(0.75), Inches(1), text_width, Inches(0.5))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "KEY CONCEPT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    concept_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), text_width, Inches(1))
    tf = concept_box.text_frame
    p = tf.paragraphs[0]
    p.text = concept
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    def_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), text_width, Inches(2))
    tf = def_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = definition
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_PRIMARY

    if example:
        ex_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.2), text_width, Inches(1.5))
        tf = ex_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Example:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        p2 = tf.add_paragraph()
        p2.text = example
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_SECONDARY
        p2.space_before = Pt(6)

    if image:
        add_image_if_exists(slide, image, Inches(8.2), Inches(1.5), width=Inches(4.8))


def add_formula_slide(prs, title, formula, explanation_lines, note=None):
    """Add a slide with a prominent formula."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Formula box with background
    formula_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(1.5), Inches(1.8),
                                         Inches(10.333), Inches(1.5))
    formula_bg.fill.solid()
    formula_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFA)  # Light blue
    formula_bg.line.color.rgb = PRIMARY

    formula_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.1), Inches(10.333), Inches(1))
    tf = formula_box.text_frame
    p = tf.paragraphs[0]
    p.text = formula
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Explanation
    explain_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.8), Inches(11.833), Inches(2.5))
    tf = explain_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(explanation_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + line
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(10)

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(11.833), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEXT_MUTED


def add_worked_example_slide(prs, title, scenario, data_rows, result, interpretation, image=None):
    """Add a slide with a worked CEA example."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Scenario
    scenario_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.833), Inches(0.8))
    tf = scenario_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = scenario
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = TEXT_SECONDARY

    # Simple data display
    data_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(6), Inches(2.5))
    tf = data_box.text_frame
    tf.word_wrap = True

    for i, row in enumerate(data_rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = row
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(6)

    # Result box
    result_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(7.5), Inches(2.2),
                                        Inches(5), Inches(1.2))
    result_bg.fill.solid()
    result_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)  # Light green
    result_bg.line.color.rgb = SUCCESS

    result_box = slide.shapes.add_textbox(Inches(7.7), Inches(2.4), Inches(4.6), Inches(1))
    tf = result_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ICER ="
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SUCCESS

    p2 = tf.add_paragraph()
    p2.text = result
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = SUCCESS

    # Interpretation
    interp_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.2), Inches(11.833), Inches(1.5))
    tf = interp_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Interpretation:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p2 = tf.add_paragraph()
    p2.text = interpretation
    p2.font.size = Pt(18)
    p2.font.color.rgb = TEXT_PRIMARY
    p2.space_before = Pt(6)

    if image:
        add_image_if_exists(slide, image, Inches(8.5), Inches(4.5), width=Inches(4))


def add_comparison_table_slide(prs, title, headers, rows, note=None):
    """Add a slide with a comparison table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Create table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    col_width = Inches(11.833 / num_cols)

    table = slide.shapes.add_table(num_rows, num_cols,
                                    Inches(0.75), Inches(1.6),
                                    Inches(11.833), Inches(4.5)).table

    # Style header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(16)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    # Fill data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14)
            para.font.color.rgb = TEXT_PRIMARY

            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xF8)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(11.833), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEXT_MUTED


def add_threshold_slide(prs):
    """Add the decision thresholds slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Decision Thresholds: What's \"Cost-Effective\"?"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Three zones
    zone_width = Inches(3.8)
    zone_height = Inches(2.5)
    zone_top = Inches(1.8)

    # Zone 1: Highly cost-effective (green)
    zone1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(0.75), zone_top,
                                    zone_width, zone_height)
    zone1.fill.solid()
    zone1.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    zone1.line.color.rgb = SUCCESS

    z1_title = slide.shapes.add_textbox(Inches(0.9), Inches(2), zone_width - Inches(0.3), Inches(0.5))
    tf = z1_title.text_frame
    p = tf.paragraphs[0]
    p.text = "< $50,000/QALY"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.alignment = PP_ALIGN.CENTER

    z1_text = slide.shapes.add_textbox(Inches(0.9), Inches(2.6), zone_width - Inches(0.3), Inches(1.5))
    tf = z1_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Highly Cost-Effective"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Generally accepted as good value. Strong case for adoption."
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.alignment = PP_ALIGN.CENTER

    # Zone 2: Potentially cost-effective (yellow)
    zone2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(4.75), zone_top,
                                    zone_width, zone_height)
    zone2.fill.solid()
    zone2.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    zone2.line.color.rgb = ACCENT

    z2_title = slide.shapes.add_textbox(Inches(4.9), Inches(2), zone_width - Inches(0.3), Inches(0.5))
    tf = z2_title.text_frame
    p = tf.paragraphs[0]
    p.text = "$50,000–$150,000/QALY"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    z2_text = slide.shapes.add_textbox(Inches(4.9), Inches(2.6), zone_width - Inches(0.3), Inches(1.5))
    tf = z2_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Potentially Cost-Effective"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Context-dependent. May be acceptable depending on budget, equity, other factors."
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.alignment = PP_ALIGN.CENTER

    # Zone 3: Not cost-effective (red/orange)
    zone3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(8.75), zone_top,
                                    zone_width, zone_height)
    zone3.fill.solid()
    zone3.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    zone3.line.color.rgb = WARNING

    z3_title = slide.shapes.add_textbox(Inches(8.9), Inches(2), zone_width - Inches(0.3), Inches(0.5))
    tf = z3_title.text_frame
    p = tf.paragraphs[0]
    p.text = "> $150,000/QALY"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WARNING
    p.alignment = PP_ALIGN.CENTER

    z3_text = slide.shapes.add_textbox(Inches(8.9), Inches(2.6), zone_width - Inches(0.3), Inches(1.5))
    tf = z3_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Less Cost-Effective"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WARNING
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Higher bar to justify. May still be appropriate for equity or other reasons."
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.alignment = PP_ALIGN.CENTER

    # Bottom note
    note_box = slide.shapes.add_textbox(Inches(0.75), Inches(5), Inches(11.833), Inches(1.5))
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Important: These thresholds are guidelines, not rules. The U.S. has no official threshold. The $50K–$150K range reflects common practice in research and policy discussions."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_SECONDARY

    p2 = tf.add_paragraph()
    p2.text = "Some argue thresholds should be higher ($200K+) based on the value of a statistical life."
    p2.font.size = Pt(14)
    p2.font.italic = True
    p2.font.color.rgb = TEXT_MUTED
    p2.space_before = Pt(8)


def add_beyond_cea_slide(prs):
    """Add slide about when to choose a program even if not cost-effective."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "When \"Not Cost-Effective\" Isn't the End of the Story"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Key message
    msg_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(11.833), Inches(0.8))
    tf = msg_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Cost-effectiveness is one input to decisions, not the only input."
    p.font.size = Pt(22)
    p.font.italic = True
    p.font.color.rgb = TEXT_SECONDARY

    # Reasons to proceed anyway
    reasons_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.4), Inches(11.833), Inches(4))
    tf = reasons_box.text_frame
    tf.word_wrap = True

    reasons = [
        ("Equity considerations", "A program serving underserved populations may be worth a higher cost per outcome."),
        ("Legal or ethical obligations", "Some services must be provided regardless of cost-effectiveness."),
        ("Political or community priorities", "Values and preferences matter beyond efficiency."),
        ("No alternatives exist", "An expensive option may beat doing nothing."),
        ("Strategic investments", "Building capacity now may pay off later in ways CEA doesn't capture."),
    ]

    for i, (title, desc) in enumerate(reasons):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {title}:"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.space_before = Pt(8) if i > 0 else Pt(0)

        p2 = tf.add_paragraph()
        p2.text = f"   {desc}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_SECONDARY
        p2.space_after = Pt(4)


def add_stay_connected_slide(prs):
    """Add Stay Connected CTA slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Stay Connected"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.833), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Continue your health economics learning journey"
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_SECONDARY

    # Three columns
    col_width = Inches(3.8)
    col_top = Inches(2.3)

    # Column 1
    col1_title = slide.shapes.add_textbox(Inches(0.75), col_top, col_width, Inches(0.5))
    tf = col1_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Practice in Methods Lab"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    col1_content = slide.shapes.add_textbox(Inches(0.75), Inches(2.9), col_width, Inches(1.5))
    tf = col1_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Interactive labs on ICER calculation, QALYs, and decision thresholds."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "caphegroup.org/methods-lab"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.space_before = Pt(12)

    # Column 2
    col2_title = slide.shapes.add_textbox(Inches(4.75), col_top, col_width, Inches(0.5))
    tf = col2_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Get Feedback"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    col2_content = slide.shapes.add_textbox(Inches(4.75), Inches(2.9), col_width, Inches(1.5))
    tf = col2_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Bring your CEA questions to monthly peer review sessions."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "info@caphegroup.org"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.space_before = Pt(12)

    # Column 3
    col3_title = slide.shapes.add_textbox(Inches(8.75), col_top, col_width, Inches(0.5))
    tf = col3_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Next Webinar"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    col3_content = slide.shapes.add_textbox(Inches(8.75), Inches(2.9), col_width, Inches(1.5))
    tf = col3_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "June 11: Return on Investment in Public Health"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "caphegroup.org/programs"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.space_before = Pt(12)

    # Methods Lab highlight box
    lab_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.75), Inches(5),
                                      Inches(11.833), Inches(1.8))
    lab_box.fill.solid()
    lab_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    lab_box.line.color.rgb = PRIMARY

    lab_title = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.333), Inches(0.4))
    tf = lab_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Related Methods Lab Modules"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    lab_content = slide.shapes.add_textbox(Inches(1), Inches(5.65), Inches(11.333), Inches(1))
    tf = lab_content.text_frame
    p = tf.paragraphs[0]
    p.text = "The Cost-Effectiveness Ratio  •  Measuring Health: QALYs  •  Decision Thresholds  •  Choosing the Comparator"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY


def add_closing_slide(prs):
    """Add closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.5))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(1.5))
    tf = contact_box.text_frame

    for i, contact in enumerate(["caphegroup.org", "info@caphegroup.org"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = contact
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)


def main():
    """Generate the complete presentation."""
    prs = create_presentation()

    print("Generating slides for: Understanding Cost-Effectiveness: The Basics")
    print("April 9, 2026 Webinar")
    print("=" * 70)

    # ==========================================================================
    # SLIDE 1: Title
    # ==========================================================================
    add_title_slide(
        prs,
        title="Understanding Cost-Effectiveness:\nThe Basics",
        subtitle="Free Webinar Series",
        date="April 9, 2026 · 12:00 PM PT"
    )
    print("  1. Title slide")

    # ==========================================================================
    # SLIDE 2: Quick Recap
    # ==========================================================================
    add_content_slide(
        prs,
        title="Quick Recap: Key Concepts from February",
        bullets=[
            "Scarcity: Resources are limited. We can't do everything",
            "Opportunity cost: Every choice means giving up something else",
            "Trade-offs: Breadth vs. depth, efficiency vs. equity",
            "Economic thinking helps us prioritize"
        ],
        note="Missed our February webinar? Recording available at caphegroup.org/programs"
    )
    print("  2. Quick recap")

    # ==========================================================================
    # SLIDE 3: Today's Objectives
    # ==========================================================================
    add_content_slide(
        prs,
        title="Today's Objectives",
        bullets=[
            "Understand what cost-effectiveness analysis (CEA) is and when to use it",
            "Learn to interpret the Incremental Cost-Effectiveness Ratio (ICER)",
            "Explore how we measure health outcomes (QALYs and alternatives)",
            "Know what decision thresholds mean. Understand their limitations",
            "See a worked example from California public health"
        ]
    )
    print("  3. Today's objectives")

    # ==========================================================================
    # SLIDE 4: Section - What is CEA?
    # ==========================================================================
    add_section_slide(prs, "Part 1", "What is Cost-Effectiveness Analysis?")
    print("  4. Section: What is CEA?")

    # ==========================================================================
    # SLIDE 5: The Question CEA Answers
    # ==========================================================================
    add_key_concept_slide(
        prs,
        concept="The Question CEA Answers",
        definition="\"Given our limited resources, which option gives us the most health improvement per dollar spent?\"",
        example="If we have $1 million, should we expand diabetes screening or fund a walking program? CEA helps compare options using a common metric.",
        image="cea-concept-question.png"
    )
    print("  5. The question CEA answers")

    # ==========================================================================
    # SLIDE 6: What CEA Is and Isn't
    # ==========================================================================
    add_content_slide(
        prs,
        title="What CEA Is (and Isn't)",
        bullets=[
            "CEA IS a tool for comparing value across interventions",
            "CEA IS a way to summarize costs and outcomes in a single ratio",
            "CEA IS NOT a decision-maker. It informs decisions, doesn't make them",
            "CEA IS NOT a judgment of worth. It measures efficiency",
            "CEA IS NOT the only consideration (equity, feasibility, politics matter too)"
        ],
        image="cea-concept-icer-ratio.png"
    )
    print("  6. What CEA is and isn't")

    # ==========================================================================
    # SLIDE 7: Section - The ICER
    # ==========================================================================
    add_section_slide(prs, "Part 2", "The ICER: The Core Metric")
    print("  7. Section: The ICER")

    # ==========================================================================
    # SLIDE 8: The ICER Formula
    # ==========================================================================
    add_formula_slide(
        prs,
        title="The Incremental Cost-Effectiveness Ratio",
        formula="ICER = (Cost_new − Cost_current) ÷ (Effect_new − Effect_current)",
        explanation_lines=[
            "Cost_new: Total cost of the new program or intervention",
            "Cost_current: Total cost of current practice (or doing nothing)",
            "Effect_new: Health outcomes achieved by the new program",
            "Effect_current: Health outcomes under current practice",
            "Result: Cost per additional unit of health gained"
        ],
        note="The ICER tells you: 'How much extra do I pay for each extra unit of health?'"
    )
    print("  8. The ICER formula")

    # ==========================================================================
    # SLIDE 9: Worked Example 1 - Simple
    # ==========================================================================
    add_worked_example_slide(
        prs,
        title="Worked Example 1: Diabetes Screening Program",
        scenario="Your county is considering expanding diabetes screening to a new population.",
        data_rows=[
            "Current practice: $500,000/year, prevents 50 cases",
            "Expanded screening: $800,000/year, prevents 80 cases",
            "─────────────────────────────────",
            "Incremental cost: $800K − $500K = $300,000",
            "Incremental effect: 80 − 50 = 30 cases prevented"
        ],
        result="$300,000 ÷ 30 = $10,000 per case prevented",
        interpretation="For every additional diabetes case we want to prevent, we'd spend $10,000. Is that worth it? Depends on downstream costs of diabetes and decision thresholds."
    )
    print("  9. Worked example 1")

    # ==========================================================================
    # SLIDE 10: Worked Example 2 - CA Tobacco
    # ==========================================================================
    add_worked_example_slide(
        prs,
        title="California Example: Tobacco Control Program",
        scenario="California's Tobacco Control Program (est. 1989) funded by Prop 99 cigarette tax.",
        data_rows=[
            "Program cost: ~$1.8 billion (1989-2004)",
            "Healthcare savings: ~$86 billion",
            "Lives saved: Estimated 1+ million life-years",
            "─────────────────────────────────",
            "ICER: ~$2,000–$4,000 per life-year saved"
        ],
        result="$2,000–$4,000 per life-year saved",
        interpretation="Extremely cost-effective by any standard. Prevention at scale works. This analysis helped justify continued funding and inspired other states.",
        image="cea-example-tobacco.png"
    )
    print("  10. CA Tobacco example")

    # ==========================================================================
    # SLIDE 11: Decision Thresholds
    # ==========================================================================
    add_threshold_slide(prs)
    print("  11. Decision thresholds")

    # ==========================================================================
    # SLIDE 12: Beyond Cost-Effectiveness
    # ==========================================================================
    add_beyond_cea_slide(prs)
    print("  12. Beyond cost-effectiveness")

    # ==========================================================================
    # SLIDE 13: Section - Measuring Outcomes
    # ==========================================================================
    add_section_slide(prs, "Part 3", "Measuring Health Outcomes")
    print("  13. Section: Measuring outcomes")

    # ==========================================================================
    # SLIDE 14: The Measurement Problem
    # ==========================================================================
    add_content_slide(
        prs,
        title="The Measurement Problem",
        bullets=[
            "Program A prevents 100 deaths",
            "Program B prevents 500 hospitalizations",
            "Program C improves quality of life for 1,000 people",
            "How do we compare them?",
            "  We need a common unit that captures both mortality and morbidity"
        ]
    )
    print("  14. The measurement problem")

    # ==========================================================================
    # SLIDE 15: Outcome Measures Comparison
    # ==========================================================================
    add_comparison_table_slide(
        prs,
        title="Ways to Measure Health Outcomes",
        headers=["Measure", "What It Captures", "Best For"],
        rows=[
            ["Lives saved", "Mortality only", "Life-or-death interventions"],
            ["Life-years gained", "Years of life added", "Comparing across age groups"],
            ["Cases prevented", "Disease incidence", "Disease-specific programs"],
            ["Hospitalizations averted", "Acute care use", "Budget impact, system capacity"],
            ["QALYs", "Years × quality (0-1)", "Comparing across conditions"],
            ["DALYs", "Years lost to death + disability", "Global health, burden of disease"],
        ],
        note="No measure is perfect. Choose based on your audience and what data you have."
    )
    print("  15. Outcome measures table")

    # ==========================================================================
    # SLIDE 16: QALYs Explained
    # ==========================================================================
    add_key_concept_slide(
        prs,
        concept="QALYs: Quality-Adjusted Life Years",
        definition="One QALY = one year of life in perfect health. A year with reduced quality of life counts as less than 1 QALY.",
        example="A treatment that extends life by 5 years but at 0.6 quality = 5 × 0.6 = 3 QALYs gained. QALYs capture both quantity and quality of life.",
        image="cea-concept-measures.png"
    )
    print("  16. QALYs explained")

    # ==========================================================================
    # SLIDE 17: QALYs - The Controversy
    # ==========================================================================
    add_content_slide(
        prs,
        title="QALYs: A Note on Controversy",
        bullets=[
            "QALYs are widely used but not without criticism",
            "Concern: May undervalue care for disabled or elderly populations",
            "  'A year of life with disability' should not automatically count less",
            "Some argue it reflects societal bias against disability",
            "Alternative: Equal value to all life-years, focus on patient preferences",
            "For today: Know QALYs exist, understand the debate, use thoughtfully"
        ],
        note="This is an active area of ethical and methodological discussion in health economics."
    )
    print("  17. QALY controversy")

    # ==========================================================================
    # SLIDE 18: Practical Guidance
    # ==========================================================================
    add_content_slide(
        prs,
        title="Which Measure Should You Use?",
        bullets=[
            "For county-level work: Natural units are often most practical",
            "  Cases prevented, hospitalizations averted, ED visits avoided",
            "  Easier to explain to supervisors and policymakers",
            "For comparing across conditions: QALYs allow apples-to-apples",
            "For published research: QALYs are standard; you'll see them cited",
            "For grant applications: Match what the funder expects",
            "When in doubt: Report multiple measures for transparency"
        ]
    )
    print("  18. Which measure to use")

    # ==========================================================================
    # SLIDE 19: Section - Limitations
    # ==========================================================================
    add_section_slide(prs, "Part 4", "Limitations and Practical Considerations")
    print("  19. Section: Limitations")

    # ==========================================================================
    # SLIDE 20: What CEA Doesn't Tell You
    # ==========================================================================
    add_content_slide(
        prs,
        title="What CEA Doesn't Tell You",
        bullets=[
            "Affordability: Something can be cost-effective but unaffordable",
            "  Budget impact analysis addresses this separately",
            "Equity: CEA maximizes total health, not distribution",
            "  May favor programs that help those already well-off",
            "Feasibility: A program may be cost-effective but impossible to implement",
            "Uncertainty: Point estimates hide ranges and assumptions",
            "Political context: Evidence isn't always persuasive in policy"
        ],
        image="cea-concept-beyond.png"
    )
    print("  20. What CEA doesn't tell you")

    # ==========================================================================
    # SLIDE 21: Common Pitfalls
    # ==========================================================================
    add_content_slide(
        prs,
        title="Common Pitfalls in CEA",
        bullets=[
            "Wrong comparator: Comparing to 'nothing' instead of current practice",
            "Narrow cost perspective: Missing costs borne by patients/families",
            "Short time horizon: Missing long-term benefits or costs",
            "Ignoring uncertainty: Treating estimates as facts",
            "Cherry-picking outcomes: Reporting only favorable measures",
            "Forgetting context: What's cost-effective in one setting may not be elsewhere"
        ]
    )
    print("  21. Common pitfalls")

    # ==========================================================================
    # SLIDE 22: Getting Started
    # ==========================================================================
    add_content_slide(
        prs,
        title="Getting Started with CEA in Your Work",
        bullets=[
            "Start by reading existing CEAs on similar interventions",
            "  CEA Registry (Tufts): cear.org is a searchable database",
            "Identify your comparator: What's the current practice?",
            "Document costs carefully: Include all relevant cost categories",
            "Choose outcomes your audience understands",
            "Acknowledge uncertainty: Report ranges, not just point estimates",
            "Partner with economists or academics if you need help"
        ]
    )
    print("  22. Getting started")

    # ==========================================================================
    # SLIDE 23: Stay Connected
    # ==========================================================================
    add_stay_connected_slide(prs)
    print("  23. Stay Connected")

    # ==========================================================================
    # SLIDE 24: Closing
    # ==========================================================================
    add_closing_slide(prs)
    print("  24. Closing slide")

    # Save presentation
    output_dir = Path("/Users/victoriaperez/Projects/CAPHE/07_website/outputs/presentations")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / "CAPHE_Webinar_CEABasics_2026-04-09.pptx"

    prs.save(str(output_path))

    print("=" * 70)
    print(f"\n✓ Presentation saved: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    return output_path


if __name__ == "__main__":
    main()
