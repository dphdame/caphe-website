#!/usr/bin/env python3
"""
Generate CAPHE-branded slides for:
"Introduction to Health Economics for Public Health"
February 12, 2026 - Free Webinar

Target audience: County health department staff, policy analysts
Duration: 1 hour (12:00 PM - 1:00 PM PT)

VERSION 2: Updated with visuals, Epi+Economics framing, Stay Connected CTA
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# =============================================================================
# CAPHE Brand Colors (matching caphegroup.org)
# =============================================================================
PRIMARY = RGBColor(0x00, 0x30, 0x80)      # Institutional Blue
PRIMARY_LIGHT = RGBColor(0x00, 0x41, 0xA5)
ACCENT = RGBColor(0xDA, 0x77, 0x0D)       # Poppy Gold
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_WARM = RGBColor(0xFA, 0xFA, 0xF8)
TEXT_PRIMARY = RGBColor(0x1C, 0x1C, 0x1C)
TEXT_SECONDARY = RGBColor(0x42, 0x42, 0x42)
TEXT_MUTED = RGBColor(0x64, 0x64, 0x64)

# Image paths
IMAGES_DIR = Path("/Users/victoriaperez/Projects/CAPHE/07_website/outputs/presentations/images")


def create_presentation():
    """Create presentation with 16:9 widescreen dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_image_if_exists(slide, image_name, left, top, width=None, height=None):
    """Add image to slide if it exists, otherwise skip silently."""
    image_path = IMAGES_DIR / image_name
    if image_path.exists():
        if width and height:
            slide.shapes.add_picture(str(image_path), left, top, width, height)
        elif width:
            slide.shapes.add_picture(str(image_path), left, top, width=width)
        elif height:
            slide.shapes.add_picture(str(image_path), left, top, height=height)
        else:
            slide.shapes.add_picture(str(image_path), left, top)
        return True
    return False


def add_title_slide(prs, title, subtitle, date):
    """Add title slide with CAPHE branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_WARM
    bg.line.fill.background()

    # Top accent bar - Institutional Blue
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY
    top_bar.line.fill.background()

    # CAPHE header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CAPHE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "California Association of Public Health Economists"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.alignment = PP_ALIGN.CENTER

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.2), Inches(11.833), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Subtitle/Presenter
    sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.8), Inches(11.833), Inches(0.75))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.8), Inches(11.833), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = date
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER

    # Bottom accent bar - Poppy Gold
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = ACCENT
    bottom_bar.line.fill.background()


def add_section_slide(prs, title, subtitle=None):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3), Inches(11.833), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = ACCENT
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)


def add_content_slide(prs, title, bullets, note=None, image=None):
    """Add content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Adjust content width if image present
    content_width = Inches(7) if image else Inches(11.833)

    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), content_width, Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle sub-bullets (indented with spaces or tabs)
        if bullet.startswith("  ") or bullet.startswith("\t"):
            p.text = "    " + bullet.strip()
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.text = "• " + bullet.strip()
            p.font.size = Pt(22)
            p.level = 0

        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(12)

    # Optional image on right side
    if image:
        add_image_if_exists(slide, image, Inches(8.5), Inches(1.6), width=Inches(4.5))

    # Optional note at bottom
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(11.833), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEXT_MUTED


def add_two_column_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    """Add slide with two columns."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(5.5), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.1), Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    # Right column content
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.1), Inches(5.5), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(8)


def add_key_concept_slide(prs, concept, definition, example=None, image=None):
    """Add a slide highlighting a key concept."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Adjust text width if image present
    text_width = Inches(7) if image else Inches(11.833)

    # Concept label
    label_box = slide.shapes.add_textbox(Inches(0.75), Inches(1), text_width, Inches(0.5))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "KEY CONCEPT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    # Concept term
    concept_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), text_width, Inches(1))
    tf = concept_box.text_frame
    p = tf.paragraphs[0]
    p.text = concept
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Definition
    def_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), text_width, Inches(2))
    tf = def_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = definition
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_PRIMARY

    # Example (if provided)
    if example:
        ex_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.2), text_width, Inches(1.5))
        tf = ex_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Example:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        p2 = tf.add_paragraph()
        p2.text = example
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_SECONDARY
        p2.space_before = Pt(6)

    # Optional image on right side
    if image:
        add_image_if_exists(slide, image, Inches(8.2), Inches(1.5), width=Inches(4.8))


def add_epi_econ_slide(prs):
    """Add the Epidemiology + Economics complementary disciplines slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Building on What You Know"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Subtitle - key message
    sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.833), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Epidemiology and economics work together to inform better decisions"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = TEXT_SECONDARY

    # Two-column comparison
    # Left column - Epidemiology
    left_title = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(5.5), Inches(0.5))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Epidemiology asks..."
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    left_content = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(5.5), Inches(2.5))
    tf = left_content.text_frame
    tf.word_wrap = True

    epi_questions = [
        "Does this intervention reduce disease?",
        "What is the effect size?",
        "Is the evidence strong?",
        "Who benefits most?"
    ]
    for i, q in enumerate(epi_questions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + q
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(8)

    # Right column - Economics
    right_title = slide.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.5), Inches(0.5))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Economics adds..."
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    right_content = slide.shapes.add_textbox(Inches(7), Inches(2.8), Inches(5.5), Inches(2.5))
    tf = right_content.text_frame
    tf.word_wrap = True

    econ_questions = [
        "Is this the best use of limited resources?",
        "What is the cost per outcome?",
        "What are we giving up to do this?",
        "How do we maximize impact?"
    ]
    for i, q in enumerate(econ_questions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + q
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(8)

    # Bottom message - collaborative framing
    bottom_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.8), Inches(11.833), Inches(1))
    tf = bottom_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "You likely already have strong skills in identifying what works. Health economics helps you make the case for why it's worth doing—and what to prioritize when you can't do everything."
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_SECONDARY
    p.alignment = PP_ALIGN.CENTER

    # Add image if available (right side, overlapping columns slightly)
    add_image_if_exists(slide, "slide-concept-epi-econ-complement.png",
                        Inches(9), Inches(2.5), width=Inches(4))


def add_stay_connected_slide(prs):
    """Add Stay Connected CTA slide before closing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Stay Connected"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.833), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Join a community of California public health economists"
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_SECONDARY

    # Three columns for CTAs
    # Column 1: Email Updates
    col1_title = slide.shapes.add_textbox(Inches(0.75), Inches(2.3), Inches(3.8), Inches(0.5))
    tf = col1_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Get Updates"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    col1_content = slide.shapes.add_textbox(Inches(0.75), Inches(2.9), Inches(3.8), Inches(1.5))
    tf = col1_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Sign up for our newsletter to hear about upcoming webinars, new Methods Lab modules, and resources."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "caphegroup.org/subscribe"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.space_before = Pt(12)

    # Column 2: Ask Questions
    col2_title = slide.shapes.add_textbox(Inches(4.75), Inches(2.3), Inches(3.8), Inches(0.5))
    tf = col2_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Ask Questions"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    col2_content = slide.shapes.add_textbox(Inches(4.75), Inches(2.9), Inches(3.8), Inches(1.5))
    tf = col2_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Have a question about applying health economics in your work? Reach out—we're here to help."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "info@caphegroup.org"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.space_before = Pt(12)

    # Column 3: Continue Learning
    col3_title = slide.shapes.add_textbox(Inches(8.75), Inches(2.3), Inches(3.8), Inches(0.5))
    tf = col3_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Continue Learning"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    col3_content = slide.shapes.add_textbox(Inches(8.75), Inches(2.9), Inches(3.8), Inches(1.5))
    tf = col3_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Explore our free Methods Lab tutorials to practice these concepts at your own pace."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "caphegroup.org/methods-lab"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.space_before = Pt(12)

    # Upcoming webinars box
    upcoming_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(0.75), Inches(5),
                                           Inches(11.833), Inches(1.8))
    upcoming_box.fill.solid()
    upcoming_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)  # Light blue-gray
    upcoming_box.line.color.rgb = PRIMARY

    upcoming_title = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.333), Inches(0.4))
    tf = upcoming_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Upcoming Webinars"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    upcoming_content = slide.shapes.add_textbox(Inches(1), Inches(5.65), Inches(11.333), Inches(1))
    tf = upcoming_content.text_frame
    p = tf.paragraphs[0]
    p.text = "April 9: Understanding Cost-Effectiveness: The Basics  •  June 11: Return on Investment in Public Health"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "Register at caphegroup.org/programs"
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_before = Pt(6)

    # Add community image if available
    add_image_if_exists(slide, "slide-cta-community.png",
                        Inches(9.5), Inches(0.3), height=Inches(1.8))


def add_closing_slide(prs):
    """Add closing slide with CAPHE info and resources."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Blue background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # Thank you
    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Questions prompt
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.5))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(1.5))
    tf = contact_box.text_frame

    contacts = [
        "caphegroup.org",
        "info@caphegroup.org"
    ]

    for i, contact in enumerate(contacts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = contact
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)  # Light blue-white
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)


def main():
    """Generate the complete presentation."""
    prs = create_presentation()

    print("Generating slides for: Introduction to Health Economics for Public Health")
    print("VERSION 2: With visuals, Epi+Econ framing, Stay Connected CTA")
    print("=" * 70)

    # ==========================================================================
    # SLIDE 1: Title
    # ==========================================================================
    add_title_slide(
        prs,
        title="Introduction to Health Economics\nfor Public Health",
        subtitle="Free Webinar Series",
        date="February 12, 2026 · 12:00 PM PT"
    )
    print("  1. Title slide")

    # ==========================================================================
    # SLIDE 2: Welcome & Objectives
    # ==========================================================================
    add_content_slide(
        prs,
        title="Today's Objectives",
        bullets=[
            "Understand what health economics is and why it matters for public health",
            "Learn key economic concepts: scarcity, trade-offs, opportunity cost",
            "Explore how economic thinking can strengthen program decisions",
            "Discover resources for continued learning"
        ],
        note="This is an introductory session—no prior economics background required."
    )
    print("  2. Objectives slide")

    # ==========================================================================
    # SLIDE 3: Epi + Economics (NEW - complementary framing)
    # ==========================================================================
    add_epi_econ_slide(prs)
    print("  3. Epi + Economics: Building on What You Know")

    # ==========================================================================
    # SLIDE 4: Section - What is Health Economics?
    # ==========================================================================
    add_section_slide(prs, "Part 1", "What is Health Economics?")
    print("  4. Section: What is Health Economics?")

    # ==========================================================================
    # SLIDE 5: Health Economics Definition
    # ==========================================================================
    add_key_concept_slide(
        prs,
        concept="Health Economics",
        definition="The study of how individuals, institutions, and society make choices about allocating scarce resources to improve health outcomes.",
        example="Should we invest $1 million in a new vaccination program or expand mental health services? Health economics provides frameworks to answer questions like this."
    )
    print("  5. Key concept: Health Economics")

    # ==========================================================================
    # SLIDE 6: Why It Matters for Public Health
    # ==========================================================================
    add_content_slide(
        prs,
        title="Why Health Economics Matters for Public Health",
        bullets=[
            "Resources are always limited—every dollar spent is a dollar not available elsewhere",
            "Public health decisions affect population-level outcomes",
            "Stakeholders (supervisors, legislators) increasingly demand evidence of value",
            "Economic analysis helps prioritize interventions with greatest impact",
            "Strengthens grant applications and program justifications"
        ]
    )
    print("  6. Why it matters")

    # ==========================================================================
    # SLIDE 7: The Fundamental Problem - Scarcity
    # ==========================================================================
    add_key_concept_slide(
        prs,
        concept="Scarcity",
        definition="The fundamental economic problem: unlimited wants and needs, but limited resources to satisfy them.",
        example="Your county has budget for 2 new positions. Do you hire an epidemiologist, a health educator, or a data analyst? You can't hire all three.",
        image="slide-concept-scarcity.png"
    )
    print("  7. Key concept: Scarcity")

    # ==========================================================================
    # SLIDE 8: Section - Core Concepts
    # ==========================================================================
    add_section_slide(prs, "Part 2", "Core Economic Concepts")
    print("  8. Section: Core Concepts")

    # ==========================================================================
    # SLIDE 9: Opportunity Cost
    # ==========================================================================
    add_key_concept_slide(
        prs,
        concept="Opportunity Cost",
        definition="The value of the next-best alternative that you give up when making a choice. It's not just about money—it's about what else you could have done.",
        example="If your team spends 6 months developing a diabetes prevention program, the opportunity cost includes other programs that couldn't be developed during that time.",
        image="slide-concept-opportunity-cost.png"
    )
    print("  9. Key concept: Opportunity Cost")

    # ==========================================================================
    # SLIDE 10: Trade-offs
    # ==========================================================================
    add_content_slide(
        prs,
        title="Trade-offs in Public Health",
        bullets=[
            "Breadth vs. Depth: Serve more people with less intensity, or fewer people with more support?",
            "Prevention vs. Treatment: Invest upstream or address immediate needs?",
            "Equity vs. Efficiency: Maximize total health gains or prioritize underserved populations?",
            "Short-term vs. Long-term: Quick wins today or sustainable impact over years?",
            "  No right answer—depends on values, context, and evidence"
        ],
        image="slide-concept-tradeoffs.png"
    )
    print("  10. Trade-offs")

    # ==========================================================================
    # SLIDE 11: Marginal Thinking
    # ==========================================================================
    add_key_concept_slide(
        prs,
        concept="Marginal Thinking",
        definition="Focus on the additional (marginal) benefit of one more unit of activity, compared to its additional cost.",
        example="Should we add a 5th community health worker to our team? The question isn't whether CHWs are valuable—it's whether the 5th CHW adds enough benefit to justify the cost.",
        image="slide-concept-marginal-thinking.png"
    )
    print("  11. Key concept: Marginal Thinking")

    # ==========================================================================
    # SLIDE 12: Section - Types of Economic Analysis
    # ==========================================================================
    add_section_slide(prs, "Part 3", "Types of Economic Analysis")
    print("  12. Section: Types of Analysis")

    # ==========================================================================
    # SLIDE 13: Overview of Analysis Types
    # ==========================================================================
    add_two_column_slide(
        prs,
        title="Economic Analysis in Public Health",
        left_title="Cost Analysis",
        left_bullets=[
            "Cost of illness studies",
            "Budget impact analysis",
            "Cost minimization",
            "Focus: What does it cost?"
        ],
        right_title="Value Analysis",
        right_bullets=[
            "Cost-effectiveness analysis (CEA)",
            "Cost-benefit analysis (CBA)",
            "Return on investment (ROI)",
            "Focus: Is it worth it?"
        ]
    )
    print("  13. Analysis types overview")

    # ==========================================================================
    # SLIDE 14: Cost-Effectiveness Analysis
    # ==========================================================================
    add_content_slide(
        prs,
        title="Cost-Effectiveness Analysis (CEA)",
        bullets=[
            "Compares costs to health outcomes (not just dollars)",
            "Results expressed as cost per outcome achieved",
            "  e.g., Cost per life saved, cost per case prevented",
            "Allows comparison across different health interventions",
            "Most common type of economic analysis in public health"
        ],
        note="We'll cover CEA in depth in our April webinar: 'Understanding Cost-Effectiveness: The Basics'",
        image="slide-concept-cost-effectiveness.png"
    )
    print("  14. Cost-effectiveness analysis")

    # ==========================================================================
    # SLIDE 15: Return on Investment
    # ==========================================================================
    add_content_slide(
        prs,
        title="Return on Investment (ROI)",
        bullets=[
            "Compares program benefits to costs, both in dollar terms",
            "ROI = (Benefits - Costs) / Costs",
            "  An ROI of 5:1 means $5 returned for every $1 invested",
            "Appealing to policymakers because it speaks their language",
            "Challenge: Converting health outcomes to dollars is controversial"
        ],
        note="Research suggests $67-88 return per dollar invested in California public health (Brown, 2014)",
        image="slide-concept-roi.png"
    )
    print("  15. Return on investment")

    # ==========================================================================
    # SLIDE 16: Section - Practical Applications
    # ==========================================================================
    add_section_slide(prs, "Part 4", "Putting It Into Practice")
    print("  16. Section: Practice")

    # ==========================================================================
    # SLIDE 17: When to Use Economic Analysis
    # ==========================================================================
    add_content_slide(
        prs,
        title="When to Use Economic Analysis",
        bullets=[
            "Comparing alternative programs or interventions",
            "Justifying new program funding to leadership",
            "Responding to budget cut proposals",
            "Writing grant applications",
            "Strategic planning and priority-setting",
            "Communicating value to external stakeholders"
        ]
    )
    print("  17. When to use")

    # ==========================================================================
    # SLIDE 18: Common Pitfalls
    # ==========================================================================
    add_content_slide(
        prs,
        title="Common Pitfalls to Avoid",
        bullets=[
            "Ignoring opportunity costs (only counting direct program costs)",
            "Comparing apples to oranges (programs with different outcomes)",
            "Overstating benefits or understating costs",
            "Ignoring uncertainty (presenting point estimates as facts)",
            "Forgetting equity considerations",
            "  Efficient ≠ Equitable—both matter in public health"
        ]
    )
    print("  18. Common pitfalls")

    # ==========================================================================
    # SLIDE 19: Getting Started
    # ==========================================================================
    add_content_slide(
        prs,
        title="Getting Started: Practical Tips",
        bullets=[
            "Start with cost documentation—know what your programs actually cost",
            "Identify your outcomes and how they're measured",
            "Look for published CEA studies of similar interventions",
            "Partner with academic institutions (UC schools, CAPHE members)",
            "Use existing tools and templates (CDC, WHO resources)",
            "Don't let 'perfect' be the enemy of 'good enough'"
        ]
    )
    print("  19. Getting started")

    # ==========================================================================
    # SLIDE 20: CAPHE Resources
    # ==========================================================================
    add_content_slide(
        prs,
        title="CAPHE Resources for You",
        bullets=[
            "Methods Lab — Free interactive tutorials on health economics methods",
            "  caphegroup.org/methods-lab",
            "Monthly Peer Review — Get feedback on your analyses from colleagues",
            "Professional Membership — Access advanced labs and working groups",
            "Webinar Series — Continued learning throughout the year",
            "  April 9: Understanding Cost-Effectiveness: The Basics",
            "  June 11: Return on Investment in Public Health"
        ]
    )
    print("  20. CAPHE resources")

    # ==========================================================================
    # SLIDE 21: Stay Connected (NEW - CTA slide)
    # ==========================================================================
    add_stay_connected_slide(prs)
    print("  21. Stay Connected (CTA)")

    # ==========================================================================
    # SLIDE 22: Closing
    # ==========================================================================
    add_closing_slide(prs)
    print("  22. Closing slide")

    # Save presentation
    output_dir = Path("/Users/victoriaperez/Projects/CAPHE/07_website/outputs/presentations")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / "CAPHE_Webinar_HealthEconomicsIntro_2026-02-12.pptx"

    prs.save(str(output_path))

    print("=" * 70)
    print(f"\n✓ Presentation saved: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    # Check for missing images
    print("\n📷 Image status:")
    images_to_check = [
        "slide-concept-scarcity.png",
        "slide-concept-opportunity-cost.png",
        "slide-concept-tradeoffs.png",
        "slide-concept-marginal-thinking.png",
        "slide-concept-epi-econ-complement.png",
        "slide-concept-cost-effectiveness.png",
        "slide-concept-roi.png",
        "slide-cta-community.png"
    ]

    for img in images_to_check:
        path = IMAGES_DIR / img
        status = "✓" if path.exists() else "✗ (missing)"
        print(f"  {status} {img}")

    return output_path


if __name__ == "__main__":
    main()
