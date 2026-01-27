#!/usr/bin/env python3
"""
CAPHE January 28, 2026 Meeting Slides
Custom deck for quarterly updates meeting.
"""

import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# =============================================================================
# CAPHE Brand Colors
# =============================================================================
PRIMARY = RGBColor(0, 48, 128)         # #003080 - Institutional Blue
ACCENT = RGBColor(218, 119, 13)        # #DA770D - Poppy Gold
BG_WARM = RGBColor(250, 250, 248)      # #FAFAF8
TEXT_PRIMARY = RGBColor(28, 28, 28)    # #1C1C1C
TEXT_SECONDARY = RGBColor(66, 66, 66)  # #424242
TEXT_MUTED = RGBColor(100, 100, 100)   # #646464
WHITE = RGBColor(255, 255, 255)

# =============================================================================
# Configuration
# =============================================================================
SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent
OUTPUT_DIR = PROJECT_DIR / "outputs" / "presentations"
MEETING_DATE = datetime(2026, 1, 28)


def create_presentation() -> Presentation:
    """Create a new presentation with CAPHE dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    return prs


def add_top_bar(slide, prs, color=None):
    """Add top accent bar."""
    if color is None:
        color = PRIMARY
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_title_slide(prs: Presentation):
    """Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_WARM
    bg.line.fill.background()

    add_top_bar(slide, prs)

    # CAPHE title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "CAPHE"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "California Association of Public Health Economists"
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Meeting info
    info = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1.5))
    tf = info.text_frame
    p = tf.paragraphs[0]
    p.text = "Member Meeting"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "January 28, 2026"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Bottom bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_agenda_slide(prs: Presentation):
    """Custom agenda for this meeting."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, prs)

    # Header
    header = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Today's Agenda"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Agenda items
    content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(5.5))
    tf = content.text_frame
    tf.word_wrap = True

    items = [
        ("1.", "New Website Demo", False),
        ("2.", "Methods Lab Overview & New Lab Ideas", False),
        ("3.", "Announcements", True),
        ("", "   • Feb 12 Webinar: Introduction to Health Economics", False),
        ("", "   • March 25: Tim Brown on Public Health ROI", False),
        ("", "   • Webinar topic ideas for fall", False),
        ("4.", "Organization Business", True),
        ("", "   • Governance docs for review", False),
        ("", "   • Letter to Erika Pan", False),
        ("5.", "Open Discussion & Next Steps", False),
    ]

    for i, (num, text, is_header) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{num}  {text}" if num else text
        p.font.size = Pt(24) if is_header or num else Pt(20)
        p.font.color.rgb = PRIMARY if is_header else TEXT_PRIMARY
        p.font.bold = is_header
        p.space_after = Pt(8)


def add_website_slide(prs: Presentation):
    """New website demo slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    header_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.5))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = PRIMARY
    header_bar.line.fill.background()

    # Section label
    label = slide.shapes.add_textbox(Inches(0.75), Inches(0.35), Inches(11.833), Inches(0.5))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "DEMO"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    # Title
    title = slide.shapes.add_textbox(Inches(0.75), Inches(0.75), Inches(11.833), Inches(0.75))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "New Website: caphegroup.org"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Features
    content = slide.shapes.add_textbox(Inches(0.75), Inches(2), Inches(5.5), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True

    features = [
        "Public Pages",
        "   • Programs & webinars",
        "   • Methods Lab (35 interactive labs!)",
        "   • Membership information",
        "",
        "Member Portal",
        "   • Meeting recordings",
        "   • Shared documents",
        "   • Peer review signup",
    ]

    for i, text in enumerate(features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        is_header = bool(text and not text.startswith("   "))
        p.font.size = Pt(22) if is_header else Pt(18)
        p.font.bold = is_header
        p.font.color.rgb = PRIMARY if is_header else TEXT_PRIMARY
        p.space_after = Pt(6)

    # URL callout
    url_box = slide.shapes.add_textbox(Inches(7), Inches(3.5), Inches(5.5), Inches(1.5))
    tf = url_box.text_frame
    p = tf.paragraphs[0]
    p.text = "caphegroup.org"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER


def add_labs_slide(prs: Presentation):
    """Methods Lab overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar
    header_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.5))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = PRIMARY
    header_bar.line.fill.background()

    label = slide.shapes.add_textbox(Inches(0.75), Inches(0.35), Inches(11.833), Inches(0.5))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "DEMO"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    title = slide.shapes.add_textbox(Inches(0.75), Inches(0.75), Inches(11.833), Inches(0.75))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Methods Lab: 35 Interactive Learning Modules"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Topics covered
    content = slide.shapes.add_textbox(Inches(0.75), Inches(1.9), Inches(6), Inches(5))
    tf = content.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Topics Covered"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    topics = [
        "Causal inference fundamentals",
        "Threats to validity",
        "Cost-effectiveness analysis",
        "California case studies",
        "Interactive simulations",
    ]

    for topic in topics:
        p = tf.add_paragraph()
        p.text = f"• {topic}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_before = Pt(8)

    # Discussion prompt
    discuss = slide.shapes.add_shape(1, Inches(7), Inches(2.2), Inches(5.5), Inches(2.5))
    discuss.fill.solid()
    discuss.fill.fore_color.rgb = RGBColor(255, 248, 235)  # Light gold
    discuss.line.color.rgb = ACCENT

    discuss_text = slide.shapes.add_textbox(Inches(7.25), Inches(2.4), Inches(5), Inches(2.2))
    tf = discuss_text.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Discussion"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p = tf.add_paragraph()
    p.text = "What topics would be most useful for new labs?"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.space_before = Pt(12)

    # URL
    url = slide.shapes.add_textbox(Inches(7), Inches(5.5), Inches(5.5), Inches(1))
    tf = url.text_frame
    p = tf.paragraphs[0]
    p.text = "caphegroup.org/methods-lab"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER


def add_announcements_slide(prs: Presentation):
    """Announcements slide - webinars and Tim Brown."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, prs)

    header = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Announcements"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Webinar announcement
    webinar_box = slide.shapes.add_shape(1, Inches(0.75), Inches(1.6), Inches(5.5), Inches(2.2))
    webinar_box.fill.solid()
    webinar_box.fill.fore_color.rgb = RGBColor(240, 245, 255)  # Light blue
    webinar_box.line.color.rgb = PRIMARY

    webinar_text = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(5), Inches(2))
    tf = webinar_text.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "WEBINAR • FEBRUARY 12"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p = tf.add_paragraph()
    p.text = "Introduction to Health Economics"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = "12:00 PM PT • Free & open to all"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_SECONDARY
    p.space_before = Pt(8)

    # Tim Brown announcement
    tim_box = slide.shapes.add_shape(1, Inches(7), Inches(1.6), Inches(5.5), Inches(2.2))
    tim_box.fill.solid()
    tim_box.fill.fore_color.rgb = RGBColor(255, 248, 235)  # Light gold
    tim_box.line.color.rgb = ACCENT

    tim_text = slide.shapes.add_textbox(Inches(7.25), Inches(1.8), Inches(5), Inches(2))
    tf = tim_text.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "MARCH 25 MEETING"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p = tf.add_paragraph()
    p.text = "Tim Brown"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_before = Pt(8)

    p = tf.add_paragraph()
    p.text = "ROI of Public Health Spending on Hospital Utilization"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY
    p.space_before = Pt(8)

    # Discussion - webinar ideas
    discuss = slide.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(11.833), Inches(2))
    tf = discuss.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Discussion: Fall 2026 Webinar Topics"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    p = tf.add_paragraph()
    p.text = "What topics should we cover? Suggestions welcome!"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_SECONDARY
    p.space_before = Pt(12)


def add_org_docs_slide(prs: Presentation):
    """Organization documents for review."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, prs)

    header = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Governance Documents for Review"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Documents list
    content = slide.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(6), Inches(4))
    tf = content.text_frame
    tf.word_wrap = True

    docs = [
        "Articles of Organization",
        "Bylaws",
        "Conflict of Interest Policy",
        "Initial Board Meeting Minutes Template",
        "501(c)(3) Filing Checklist",
    ]

    for i, doc in enumerate(docs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {doc}"
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(12)

    # Action box
    action_box = slide.shapes.add_shape(1, Inches(7), Inches(1.7), Inches(5.5), Inches(2.5))
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = RGBColor(255, 248, 235)
    action_box.line.color.rgb = ACCENT

    action_text = slide.shapes.add_textbox(Inches(7.25), Inches(1.9), Inches(5), Inches(2.2))
    tf = action_text.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "ACTION REQUESTED"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p = tf.add_paragraph()
    p.text = "Review documents and provide feedback"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.space_before = Pt(12)

    p = tf.add_paragraph()
    p.text = "Due: February 25, 2026"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_before = Pt(12)

    # Location
    loc = slide.shapes.add_textbox(Inches(0.75), Inches(5.5), Inches(11.833), Inches(1))
    tf = loc.text_frame
    p = tf.paragraphs[0]
    p.text = "Location: Google Drive → CAPHE folder (link will be shared via email)"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_SECONDARY


def add_letter_slide(prs: Presentation):
    """Letter to Erika Pan slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, prs)

    header = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Letter to Dr. Erika Pan"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Context
    content = slide.shapes.add_textbox(Inches(0.75), Inches(1.7), Inches(6), Inches(4))
    tf = content.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "California State Epidemiologist"
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_SECONDARY

    p = tf.add_paragraph()
    p.text = ""
    p.space_before = Pt(20)

    p = tf.add_paragraph()
    p.text = "Purpose:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    points = [
        "Introduce CAPHE and our mission",
        "Highlight value for CDPH staff",
        "Explore collaboration opportunities",
    ]

    for point in points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_before = Pt(8)

    # Action box
    action_box = slide.shapes.add_shape(1, Inches(7), Inches(1.7), Inches(5.5), Inches(2.5))
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = RGBColor(255, 248, 235)
    action_box.line.color.rgb = ACCENT

    action_text = slide.shapes.add_textbox(Inches(7.25), Inches(1.9), Inches(5), Inches(2.2))
    tf = action_text.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "ACTION REQUESTED"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    p = tf.add_paragraph()
    p.text = "Review draft and provide comments"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_PRIMARY
    p.space_before = Pt(12)

    p = tf.add_paragraph()
    p.text = "Due: February 14, 2026"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.space_before = Pt(12)

    # Location
    loc = slide.shapes.add_textbox(Inches(0.75), Inches(5.5), Inches(11.833), Inches(1))
    tf = loc.text_frame
    p = tf.paragraphs[0]
    p.text = "Location: Google Drive → CAPHE → CAPHE_Letter.gdoc"
    p.font.size = Pt(18)
    p.font.color.rgb = TEXT_SECONDARY


def add_upcoming_events_slide(prs: Presentation):
    """Upcoming events slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_bar(slide, prs)

    header = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Upcoming Events"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    events = [
        ("FEB", "12", "Webinar (Free)", "Introduction to Health Economics"),
        ("FEB", "25", "Peer Review", "Monthly Member Meeting"),
        ("MAR", "25", "Peer Review", "Tim Brown: Public Health ROI"),
        ("APR", "09", "Webinar (Free)", "Cost-Effectiveness: The Basics"),
    ]

    start_y = 1.6
    for i, (month, day, event_type, title) in enumerate(events):
        # Date box
        date_box = slide.shapes.add_shape(
            1, Inches(0.75), Inches(start_y + i * 1.3), Inches(1.3), Inches(1)
        )
        date_box.fill.solid()
        date_box.fill.fore_color.rgb = PRIMARY
        date_box.line.fill.background()

        date_text = slide.shapes.add_textbox(
            Inches(0.75), Inches(start_y + i * 1.3 + 0.1), Inches(1.3), Inches(0.8)
        )
        tf = date_text.text_frame

        p = tf.paragraphs[0]
        p.text = month
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = day
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Event details
        details = slide.shapes.add_textbox(
            Inches(2.3), Inches(start_y + i * 1.3), Inches(10), Inches(1)
        )
        tf = details.text_frame

        p = tf.paragraphs[0]
        p.text = event_type
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_PRIMARY


def add_closing_slide(prs: Presentation):
    """Closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # Thank you
    thanks = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1))
    tf = thanks.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Next meeting
    next_mtg = slide.shapes.add_textbox(Inches(0.5), Inches(3.75), Inches(12.333), Inches(0.75))
    tf = next_mtg.text_frame
    p = tf.paragraphs[0]
    p.text = "Next Meeting: February 25, 2026"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 215, 240)
    p.alignment = PP_ALIGN.CENTER

    # Website
    web = slide.shapes.add_textbox(Inches(0.5), Inches(4.75), Inches(12.333), Inches(0.75))
    tf = web.text_frame
    p = tf.paragraphs[0]
    p.text = "caphegroup.org"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Email
    email = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
    tf = email.text_frame
    p = tf.paragraphs[0]
    p.text = "info@caphegroup.org"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(200, 215, 240)
    p.alignment = PP_ALIGN.CENTER


def main():
    print("=" * 60)
    print("CAPHE January 28, 2026 Meeting Slides")
    print("=" * 60)

    prs = create_presentation()

    print("\nGenerating slides...")

    add_title_slide(prs)
    print("  - Title slide")

    add_agenda_slide(prs)
    print("  - Agenda slide")

    add_website_slide(prs)
    print("  - New website demo slide")

    add_labs_slide(prs)
    print("  - Methods Lab slide")

    add_announcements_slide(prs)
    print("  - Announcements slide (webinar + Tim Brown)")

    add_org_docs_slide(prs)
    print("  - Governance documents slide")

    add_letter_slide(prs)
    print("  - Erika Pan letter slide")

    add_upcoming_events_slide(prs)
    print("  - Upcoming events slide")

    add_closing_slide(prs)
    print("  - Closing slide")

    # Save
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_path = OUTPUT_DIR / "CAPHE_Meeting_2026-01-28.pptx"
    prs.save(str(output_path))

    print(f"\n✓ Presentation saved: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    return output_path


if __name__ == "__main__":
    main()
