#!/usr/bin/env python3
"""
CAPHE Meeting Slide Generator
Automatically generates PowerPoint slides for general member meetings.

Features:
- Pulls presentation signups from Supabase
- Reads upcoming events from events.json
- Creates clean, branded slides matching caphegroup.org

Usage:
    python generate_meeting_slides.py [--date YYYY-MM-DD]

    If no date specified, generates for next peer review session.
"""

import argparse
import json
import os
import sys
from datetime import datetime, timedelta
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Optional: Supabase client (install with: pip install supabase)
try:
    from supabase import create_client, Client
    SUPABASE_AVAILABLE = True
except ImportError:
    SUPABASE_AVAILABLE = False
    print("Note: supabase package not installed. Run: pip install supabase")

# =============================================================================
# CAPHE Brand Colors (matching caphegroup.org)
# Institutional Blue + Poppy Gold - two-color palette
# =============================================================================
PRIMARY = RGBColor(0, 48, 128)         # #003080 - Institutional Blue (darker)
PRIMARY_LIGHT = RGBColor(0, 65, 165)   # #0041A5 - lighter blue
ACCENT = RGBColor(218, 119, 13)        # #DA770D - Poppy Gold (darker for readability)
ACCENT_LIGHT = RGBColor(245, 127, 23)  # #F57F17
BG_WHITE = RGBColor(255, 255, 255)     # #FFFFFF
BG_WARM = RGBColor(250, 250, 248)      # #FAFAF8
TEXT_PRIMARY = RGBColor(28, 28, 28)    # #1C1C1C
TEXT_SECONDARY = RGBColor(66, 66, 66)  # #424242
TEXT_MUTED = RGBColor(100, 100, 100)   # #646464
WHITE = RGBColor(255, 255, 255)

# =============================================================================
# Configuration
# =============================================================================
SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent
DATA_DIR = PROJECT_DIR / "data"
EVENTS_FILE = DATA_DIR / "events.json"
OUTPUT_DIR = PROJECT_DIR / "outputs" / "presentations"

# Supabase config (public keys - safe to include)
SUPABASE_URL = "https://yyetprjdxwunhtighnrq.supabase.co"
SUPABASE_ANON_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6Inl5ZXRwcmpkeHd1bmh0aWdobnJxIiwicm9sZSI6ImFub24iLCJpYXQiOjE3NjcwMzk2MDAsImV4cCI6MjA4MjYxNTYwMH0.xWguR4nFUGAflIy3iolYHUZFAY2ec0CGcFG2f8a-TWQ"


# =============================================================================
# Data Loading Functions
# =============================================================================

def load_events() -> list:
    """Load events from JSON file."""
    if not EVENTS_FILE.exists():
        print(f"Warning: Events file not found at {EVENTS_FILE}")
        return []

    with open(EVENTS_FILE) as f:
        data = json.load(f)
    return data.get("events", [])


def get_upcoming_events(events: list, from_date: datetime, count: int = 5) -> list:
    """Filter events to get upcoming ones from a given date."""
    upcoming = []
    for event in events:
        event_date = datetime.strptime(event["date"], "%Y-%m-%d")
        if event_date >= from_date:
            upcoming.append({
                **event,
                "date_obj": event_date
            })

    # Sort by date and limit
    upcoming.sort(key=lambda x: x["date_obj"])
    return upcoming[:count]


def get_next_peer_review_date(events: list, from_date: datetime = None) -> datetime:
    """Find the next peer review session date."""
    if from_date is None:
        from_date = datetime.now()

    for event in events:
        event_date = datetime.strptime(event["date"], "%Y-%m-%d")
        if event_date >= from_date and event["type"] == "Peer Review":
            return event_date

    # Fallback: last Wednesday of next month
    return from_date


def fetch_presentation_signups(meeting_date: datetime) -> list:
    """Fetch presentation signups from Supabase for a specific meeting date."""
    if not SUPABASE_AVAILABLE:
        print("Supabase not available. Skipping signup fetch.")
        return []

    try:
        supabase: Client = create_client(SUPABASE_URL, SUPABASE_ANON_KEY)

        # Format date for query
        date_str = meeting_date.strftime("%Y-%m-%d")

        # Query peer_review_requests table
        response = supabase.table("peer_review_requests").select(
            "id, topic, description, slots_requested, status, created_at"
        ).or_(
            f"meeting_date.eq.{date_str},meeting_date_2.eq.{date_str}"
        ).eq("status", "confirmed").execute()

        return response.data if response.data else []

    except Exception as e:
        print(f"Error fetching signups: {e}")
        return []


def fetch_presenter_profiles(user_ids: list) -> dict:
    """Fetch presenter profiles from Supabase."""
    if not SUPABASE_AVAILABLE or not user_ids:
        return {}

    try:
        supabase: Client = create_client(SUPABASE_URL, SUPABASE_ANON_KEY)

        response = supabase.table("profiles").select(
            "id, full_name, organization"
        ).in_("id", user_ids).execute()

        return {p["id"]: p for p in response.data} if response.data else {}

    except Exception as e:
        print(f"Error fetching profiles: {e}")
        return {}


# =============================================================================
# Slide Creation Functions
# =============================================================================

def create_presentation() -> Presentation:
    """Create a new presentation with CAPHE dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs: Presentation, meeting_date: datetime, meeting_type: str = "Member Meeting"):
    """Add title slide with meeting date and CAPHE branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background color - warm white
    background = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = BG_WARM
    background.line.fill.background()

    # Top accent bar - Institutional Blue
    top_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY
    top_bar.line.fill.background()

    # CAPHE logo/title area
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(12.333), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True

    # Main title - Institutional Blue
    p = tf.paragraphs[0]
    p.text = "CAPHE"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    p = tf.add_paragraph()
    p.text = "California Association of Public Health Economists"
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(10)

    # Meeting type and date
    date_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5), Inches(12.333), Inches(1.5)
    )
    tf = date_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = meeting_type
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT  # Poppy Gold
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = meeting_date.strftime("%B %d, %Y")
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(15)

    # Bottom accent bar - Poppy Gold
    accent_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(7.35), prs.slide_width, Inches(0.15)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()


def add_agenda_slide(prs: Presentation, presenters: list, has_events: bool = True):
    """Add agenda slide listing meeting topics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    top_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY
    top_bar.line.fill.background()

    # Header
    header = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.5), Inches(11.833), Inches(1)
    )
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Today's Agenda"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Agenda items
    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.75), Inches(11.333), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    agenda_items = ["Welcome & Announcements"]

    if presenters:
        for presenter in presenters:
            topic = presenter.get("topic", "Research Presentation")
            agenda_items.append(f"Presentation: {topic}")
    else:
        agenda_items.append("Open Discussion / Networking")

    if has_events:
        agenda_items.append("Upcoming Events")

    agenda_items.append("Q&A / Wrap-up")

    for i, item in enumerate(agenda_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i + 1}.  {item}"
        p.font.size = Pt(28)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(20)

        # Highlight presentations in Poppy Gold
        if "Presentation:" in item:
            p.font.color.rgb = ACCENT
            p.font.bold = True


def add_presenter_slide(prs: Presentation, presenter: dict, profile: dict = None):
    """Add a slide for a presenter."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Header bar - Institutional Blue
    header_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(1.5)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = PRIMARY
    header_bar.line.fill.background()

    # Presentation label - Poppy Gold
    label_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.35), Inches(11.833), Inches(0.5)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PRESENTATION"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    # Topic title
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.75), Inches(11.833), Inches(0.75)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = presenter.get("topic", "Research Presentation")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Presenter info
    presenter_name = "Presenter"
    organization = ""

    if profile:
        presenter_name = profile.get("full_name", "Presenter")
        organization = profile.get("organization", "")

    info_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2), Inches(11.833), Inches(1)
    )
    tf = info_box.text_frame

    p = tf.paragraphs[0]
    p.text = presenter_name
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY  # Institutional Blue

    if organization:
        p = tf.add_paragraph()
        p.text = organization
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_SECONDARY
        p.space_before = Pt(5)

    # Description
    description = presenter.get("description", "")
    if description:
        desc_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(3.5), Inches(11.833), Inches(3)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_PRIMARY

    # Time slots
    slots = presenter.get("slots_requested", 1)
    time_label = f"{slots * 15} minutes" if slots > 1 else "15 minutes"

    time_box = slide.shapes.add_textbox(
        Inches(10), Inches(6.5), Inches(3), Inches(0.5)
    )
    tf = time_box.text_frame
    p = tf.paragraphs[0]
    p.text = time_label
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_events_slide(prs: Presentation, events: list):
    """Add upcoming events slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    top_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, Inches(0.15)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY
    top_bar.line.fill.background()

    # Header
    header = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.5), Inches(11.833), Inches(1)
    )
    tf = header.text_frame
    p = tf.paragraphs[0]
    p.text = "Upcoming Events"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Events grid
    start_y = 1.75
    for i, event in enumerate(events):
        event_date = event.get("date_obj", datetime.now())

        # Date box - Institutional Blue
        date_box = slide.shapes.add_shape(
            1, Inches(0.75), Inches(start_y + i * 1.1), Inches(1.5), Inches(0.9)
        )
        date_box.fill.solid()
        date_box.fill.fore_color.rgb = PRIMARY
        date_box.line.fill.background()

        # Date text
        date_text = slide.shapes.add_textbox(
            Inches(0.75), Inches(start_y + i * 1.1 + 0.1), Inches(1.5), Inches(0.7)
        )
        tf = date_text.text_frame
        tf.word_wrap = False

        p = tf.paragraphs[0]
        p.text = event_date.strftime("%b").upper()
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = event_date.strftime("%d")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Event details
        details_box = slide.shapes.add_textbox(
            Inches(2.5), Inches(start_y + i * 1.1), Inches(9.833), Inches(0.9)
        )
        tf = details_box.text_frame
        tf.word_wrap = True

        # Event type badge
        event_type = event.get("type", "Event")
        is_free = not event.get("members_only", False)
        badge = f"{event_type}"
        if is_free:
            badge += " (Free)"

        p = tf.paragraphs[0]
        p.text = badge
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        p = tf.add_paragraph()
        p.text = event.get("title", "Event")
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_PRIMARY

        p = tf.add_paragraph()
        p.text = event.get("time", "TBD")
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_SECONDARY


def add_closing_slide(prs: Presentation):
    """Add closing slide with CAPHE info."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background - Institutional Blue
    background = slide.shapes.add_shape(
        1, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY
    background.line.fill.background()

    # Thank you
    thanks_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.333), Inches(1)
    )
    tf = thanks_box.text_frame

    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Website - Poppy Gold
    web_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4), Inches(12.333), Inches(0.75)
    )
    tf = web_box.text_frame

    p = tf.paragraphs[0]
    p.text = "caphegroup.org"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Contact
    contact_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5), Inches(12.333), Inches(0.5)
    )
    tf = contact_box.text_frame

    p = tf.paragraphs[0]
    p.text = "info@caphegroup.org"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(200, 215, 240)  # Light blue-white
    p.alignment = PP_ALIGN.CENTER


# =============================================================================
# Main Generation Function
# =============================================================================

def generate_meeting_slides(meeting_date: datetime = None, output_path: str = None):
    """Generate complete meeting slide deck."""

    # Load events
    events = load_events()

    # Determine meeting date
    if meeting_date is None:
        meeting_date = get_next_peer_review_date(events)
        print(f"No date specified. Using next peer review: {meeting_date.strftime('%Y-%m-%d')}")

    # Fetch presentation signups
    presenters = fetch_presentation_signups(meeting_date)
    print(f"Found {len(presenters)} confirmed presenter(s)")

    # Fetch presenter profiles
    user_ids = [p.get("user_id") for p in presenters if p.get("user_id")]
    profiles = fetch_presenter_profiles(user_ids)

    # Get upcoming events (from meeting date forward)
    upcoming_events = get_upcoming_events(events, meeting_date, count=5)
    print(f"Found {len(upcoming_events)} upcoming event(s)")

    # Create presentation
    prs = create_presentation()

    print("\nGenerating slides...")

    # 1. Title slide
    add_title_slide(prs, meeting_date, "Member Meeting")
    print("  - Title slide")

    # 2. Agenda slide
    add_agenda_slide(prs, presenters, has_events=len(upcoming_events) > 0)
    print("  - Agenda slide")

    # 3. Presenter slides (if any)
    for presenter in presenters:
        profile = profiles.get(presenter.get("user_id"), {})
        add_presenter_slide(prs, presenter, profile)
        print(f"  - Presenter: {presenter.get('topic', 'TBD')}")

    # 4. Events slide
    if upcoming_events:
        add_events_slide(prs, upcoming_events)
        print("  - Upcoming events slide")

    # 5. Closing slide
    add_closing_slide(prs)
    print("  - Closing slide")

    # Save
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    if output_path is None:
        filename = f"CAPHE_Meeting_{meeting_date.strftime('%Y-%m-%d')}.pptx"
        output_path = OUTPUT_DIR / filename
    else:
        output_path = Path(output_path)

    prs.save(str(output_path))

    print(f"\n  Presentation saved: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    return output_path


# =============================================================================
# CLI Interface
# =============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Generate CAPHE meeting slides",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
    python generate_meeting_slides.py
    python generate_meeting_slides.py --date 2026-01-28
    python generate_meeting_slides.py --output ~/Desktop/meeting.pptx
        """
    )

    parser.add_argument(
        "--date", "-d",
        help="Meeting date (YYYY-MM-DD). Default: next peer review session",
        default=None
    )

    parser.add_argument(
        "--output", "-o",
        help="Output file path. Default: outputs/presentations/CAPHE_Meeting_DATE.pptx",
        default=None
    )

    args = parser.parse_args()

    # Parse date if provided
    meeting_date = None
    if args.date:
        try:
            meeting_date = datetime.strptime(args.date, "%Y-%m-%d")
        except ValueError:
            print(f"Error: Invalid date format '{args.date}'. Use YYYY-MM-DD.")
            sys.exit(1)

    print("=" * 60)
    print("CAPHE Meeting Slide Generator")
    print("=" * 60)

    output_path = generate_meeting_slides(meeting_date, args.output)

    print("\nDone!")
    return 0


if __name__ == "__main__":
    sys.exit(main())
