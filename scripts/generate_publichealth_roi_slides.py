#!/usr/bin/env python3
"""
Generate CAPHE-branded slides for Public Health ROI Project presentation.
Covers the setup, innovation, and findings from the PublicHealthValue project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# CAPHE Brand Colors
PRIMARY = RGBColor(0x00, 0x30, 0x80)  # Dark blue
ACCENT = RGBColor(0xDA, 0x77, 0x0D)   # Orange
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)


def add_title_slide(prs, title, subtitle):
    """Add a title slide with CAPHE branding."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Blue header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(2.5))
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # CAPHE logo text
    logo_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(0.5))
    tf = logo_box.text_frame
    p = tf.paragraphs[0]
    p.text = "California Association of Public Health Economists"
    p.font.size = Pt(14)
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullets, note=None):
    """Add a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Orange accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle indented sub-bullets
        if bullet.startswith("   "):
            p.text = "    " + bullet.strip()
            p.font.size = Pt(16)
            p.level = 1
        else:
            p.text = "• " + bullet
            p.font.size = Pt(18)
            p.level = 0

        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    # Optional note at bottom
    if note:
        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    return slide


def add_findings_slide(prs, title, findings):
    """Add a findings slide with key statistics."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Orange accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Create finding boxes
    y_pos = 1.4
    for label, value, detail in findings:
        # Value box (large number)
        val_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(2.5), Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        # Label and detail
        label_box = slide.shapes.add_textbox(Inches(3.2), Inches(y_pos), Inches(6), Inches(0.8))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = PRIMARY

        p2 = tf.add_paragraph()
        p2.text = detail
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_GRAY

        y_pos += 1.1

    return slide


def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(
        prs,
        "The ROI of Public Health Spending",
        "Background for Tim Brown's March Presentation"
    )

    # Slide 2: The Brown (2014) Study
    add_content_slide(
        prs,
        "Brown (2014): Key Findings",
        [
            "Research question: What's the return on public health investment?",
            "Data: California counties, 1993-2006",
            "Method: Instrumental variables using political party control",
            "Finding: $67-88 return per dollar of public health spending",
            "   Widely cited in advocacy and budget discussions",
            "   Used to justify public health funding increases"
        ]
    )

    # Slide 3: Why This Matters
    add_content_slide(
        prs,
        "Why This Matters for California",
        [
            "Public health departments face constant budget pressure",
            "   Need credible evidence to justify funding requests",
            "ROI estimates are powerful advocacy tools",
            "   But only if the methodology is sound",
            "Counties want to demonstrate value to supervisors",
            "   Local data is more compelling than statewide averages"
        ]
    )

    # Slide 4: The Measurement Challenge
    add_content_slide(
        prs,
        "The Measurement Challenge",
        [
            "What counts as 'public health' spending?",
            "   County health budgets include many activities",
            "   Clinical care, administration, emergency response",
            "Only 40-60% is true prevention spending",
            "   Environmental health, communicable disease, maternal/child health",
            "Lumping everything together may understate ROI",
            "   Prevention dollars are doing the heavy lifting"
        ]
    )

    # Slide 5: Key Insight
    add_findings_slide(
        prs,
        "Key Insight",
        [
            ("40-60%", "True Prevention", "Share of health budgets that is actual prevention spending"),
            ("$67-88", "Brown's ROI", "Return per dollar using total health spending"),
            ("Higher?", "Refined ROI", "What if we measured only prevention spending?")
        ]
    )

    # Slide 6: Opportunity for Counties
    add_content_slide(
        prs,
        "An Opportunity for Counties",
        [
            "We want to produce county-specific ROI estimates",
            "   More useful than statewide averages for local advocacy",
            "This requires detailed spending data from counties",
            "   Breakdown by program area and function",
            "   Allows proper classification of prevention vs. other spending"
        ]
    )

    # Slide 7: What We're Asking
    add_content_slide(
        prs,
        "What We're Asking",
        [
            "Share your county's detailed expenditure data",
            "   By program area (environmental, communicable disease, etc.)",
            "   Multiple years if available",
            "We will produce a report including:",
            "   County-specific ROI estimates",
            "   Comparison across participating counties",
            "   Materials for Board of Supervisors presentations"
        ]
    )

    # Slide 8: Coming Up
    add_content_slide(
        prs,
        "Coming Up: March 25 Meeting",
        [
            "Tim Brown will present on ROI methodology",
            "   Author of the original 2014 study",
            "   'ROI of Public Health Spending on Hospital Utilization'",
            "Opportunity to discuss:",
            "   Methodological considerations",
            "   How to apply this to your county",
            "   Interest in participating in our county-level analysis"
        ]
    )

    # Save presentation
    output_dir = Path("/Users/victoriaperez/Projects/CAPHE/07_website/outputs/presentations")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / "PublicHealth_ROI_Overview.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

    return output_path


if __name__ == "__main__":
    main()
