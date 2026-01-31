#!/usr/bin/env python3
"""
Generate CAPHE-branded slides for:
"Return on Investment in Public Health"
June 11, 2026 - Free Webinar

Target audience: County health department staff, policy analysts,
                 budget officers, elected officials
Duration: 1 hour (12:00 PM - 1:00 PM PT)
Key message: "Public health is an investment, not just an expense"
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# =============================================================================
# CAPHE Brand Colors (matching caphegroup.org)
# =============================================================================
PRIMARY = RGBColor(0x00, 0x30, 0x80)      # Institutional Blue
PRIMARY_LIGHT = RGBColor(0x00, 0x41, 0xA5)
ACCENT = RGBColor(0xDA, 0x77, 0x0D)       # Poppy Gold
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_WARM = RGBColor(0xFA, 0xFA, 0xF8)
TEXT_PRIMARY = RGBColor(0x1C, 0x1C, 0x1C)
TEXT_SECONDARY = RGBColor(0x42, 0x42, 0x42)
TEXT_MUTED = RGBColor(0x64, 0x64, 0x64)
SUCCESS = RGBColor(0x2E, 0x7D, 0x32)       # Green for positive ROI
WARNING = RGBColor(0xF5, 0x7C, 0x00)       # Orange

# Image paths
IMAGES_DIR = Path("/Users/victoriaperez/Projects/CAPHE/07_website/outputs/presentations/images")


def create_presentation():
    """Create presentation with 16:9 widescreen dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_image_if_exists(slide, image_name, left, top, width=None, height=None):
    """Add image to slide if it exists, otherwise skip silently."""
    image_path = IMAGES_DIR / image_name
    if image_path.exists():
        if width and height:
            slide.shapes.add_picture(str(image_path), left, top, width, height)
        elif width:
            slide.shapes.add_picture(str(image_path), left, top, width=width)
        elif height:
            slide.shapes.add_picture(str(image_path), left, top, height=height)
        else:
            slide.shapes.add_picture(str(image_path), left, top)
        return True
    return False


def add_title_slide(prs, title, subtitle, date):
    """Add title slide with CAPHE branding."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_WARM
    bg.line.fill.background()

    # Top accent bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PRIMARY
    top_bar.line.fill.background()

    # CAPHE header
    header_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1))
    tf = header_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CAPHE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "California Association of Public Health Economists"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEXT_SECONDARY
    p2.alignment = PP_ALIGN.CENTER

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.2), Inches(11.833), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.8), Inches(11.833), Inches(0.75))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.8), Inches(11.833), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = date
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER

    # Bottom accent bar
    bottom_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
    bottom_bar.fill.solid()
    bottom_bar.fill.fore_color.rgb = ACCENT
    bottom_bar.line.fill.background()


def add_section_slide(prs, title, subtitle=None):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(3), Inches(11.833), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = ACCENT
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)


def add_content_slide(prs, title, bullets, note=None, image=None):
    """Add content slide with title and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    content_width = Inches(7) if image else Inches(11.833)

    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.6), content_width, Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if bullet.startswith("  ") or bullet.startswith("\t"):
            p.text = "    " + bullet.strip()
            p.font.size = Pt(18)
            p.level = 1
        else:
            p.text = "• " + bullet.strip()
            p.font.size = Pt(22)
            p.level = 0

        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(12)

    if image:
        add_image_if_exists(slide, image, Inches(8.5), Inches(1.6), width=Inches(4.5))

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(11.833), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEXT_MUTED


def add_key_concept_slide(prs, concept, definition, example=None, image=None):
    """Add a slide highlighting a key concept."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    text_width = Inches(7) if image else Inches(11.833)

    label_box = slide.shapes.add_textbox(Inches(0.75), Inches(1), text_width, Inches(0.5))
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    p.text = "KEY CONCEPT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    concept_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), text_width, Inches(1))
    tf = concept_box.text_frame
    p = tf.paragraphs[0]
    p.text = concept
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    def_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), text_width, Inches(2))
    tf = def_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = definition
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_PRIMARY

    if example:
        ex_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.2), text_width, Inches(1.5))
        tf = ex_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Example:"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT

        p2 = tf.add_paragraph()
        p2.text = example
        p2.font.size = Pt(18)
        p2.font.color.rgb = TEXT_SECONDARY
        p2.space_before = Pt(6)

    if image:
        add_image_if_exists(slide, image, Inches(8.2), Inches(1.5), width=Inches(4.8))


def add_formula_slide(prs, title, formula, explanation_lines, note=None):
    """Add a slide with a prominent formula."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Formula box with background
    formula_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(1.5), Inches(1.8),
                                         Inches(10.333), Inches(1.5))
    formula_bg.fill.solid()
    formula_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)  # Light green for ROI
    formula_bg.line.color.rgb = SUCCESS

    formula_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.1), Inches(10.333), Inches(1))
    tf = formula_box.text_frame
    p = tf.paragraphs[0]
    p.text = formula
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = SUCCESS
    p.alignment = PP_ALIGN.CENTER

    # Explanation
    explain_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.8), Inches(11.833), Inches(2.5))
    tf = explain_box.text_frame
    tf.word_wrap = True

    for i, line in enumerate(explanation_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + line
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(10)

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(11.833), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEXT_MUTED


def add_roi_example_slide(prs, title, program, investment, returns, roi_ratio, source=None, image=None):
    """Add a slide with an ROI example."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Program name
    program_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.833), Inches(0.5))
    tf = program_box.text_frame
    p = tf.paragraphs[0]
    p.text = program
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = TEXT_SECONDARY

    # Investment side
    inv_label = slide.shapes.add_textbox(Inches(0.75), Inches(2), Inches(5.5), Inches(0.4))
    tf = inv_label.text_frame
    p = tf.paragraphs[0]
    p.text = "INVESTMENT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    inv_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.75), Inches(2.5),
                                      Inches(5.5), Inches(1.2))
    inv_box.fill.solid()
    inv_box.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)  # Light blue
    inv_box.line.color.rgb = PRIMARY

    inv_text = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(5), Inches(1))
    tf = inv_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = investment
    p.font.size = Pt(20)
    p.font.color.rgb = PRIMARY

    # Returns side
    ret_label = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5.5), Inches(0.4))
    tf = ret_label.text_frame
    p = tf.paragraphs[0]
    p.text = "RETURNS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = SUCCESS

    ret_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(7), Inches(2.5),
                                      Inches(5.5), Inches(1.8))
    ret_box.fill.solid()
    ret_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)  # Light green
    ret_box.line.color.rgb = SUCCESS

    ret_text = slide.shapes.add_textbox(Inches(7.25), Inches(2.7), Inches(5), Inches(1.5))
    tf = ret_text.text_frame
    tf.word_wrap = True
    for i, ret in enumerate(returns):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + ret
        p.font.size = Pt(18)
        p.font.color.rgb = SUCCESS

    # ROI result
    roi_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(3.5), Inches(4.8),
                                     Inches(6.333), Inches(1.2))
    roi_bg.fill.solid()
    roi_bg.fill.fore_color.rgb = ACCENT
    roi_bg.line.fill.background()

    roi_text = slide.shapes.add_textbox(Inches(3.5), Inches(5), Inches(6.333), Inches(1))
    tf = roi_text.text_frame
    p = tf.paragraphs[0]
    p.text = f"ROI: {roi_ratio}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if source:
        source_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(11.833), Inches(0.5))
        tf = source_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Source: {source}"
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = TEXT_MUTED

    if image:
        add_image_if_exists(slide, image, Inches(0.75), Inches(4.5), width=Inches(2.5))


def add_comparison_table_slide(prs, title, headers, rows, note=None):
    """Add a slide with a comparison table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Create table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    col_width = Inches(11.833 / num_cols)

    table = slide.shapes.add_table(num_rows, num_cols,
                                    Inches(0.75), Inches(1.6),
                                    Inches(11.833), Inches(4.5)).table

    # Style header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(16)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    # Fill data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14)
            para.font.color.rgb = TEXT_PRIMARY

            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xF8)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(11.833), Inches(0.5))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(14)
        p.font.italic = True
        p.font.color.rgb = TEXT_MUTED


def add_audience_framing_slide(prs):
    """Add slide about speaking to different audiences."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Speaking to Your Audience"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    # Two columns
    col_width = Inches(5.8)
    col_top = Inches(1.6)

    # Health audience column
    health_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.75), col_top,
                                        col_width, Inches(4.5))
    health_bg.fill.solid()
    health_bg.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    health_bg.line.color.rgb = PRIMARY

    health_title = slide.shapes.add_textbox(Inches(1), Inches(1.8), col_width - Inches(0.5), Inches(0.5))
    tf = health_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Health Leaders"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    health_content = slide.shapes.add_textbox(Inches(1), Inches(2.4), col_width - Inches(0.5), Inches(3.5))
    tf = health_content.text_frame
    tf.word_wrap = True

    health_points = [
        "Focus on health outcomes first",
        "ROI supports the mission",
        "Use alongside CEA and equity",
        "\"This investment saves lives AND money\""
    ]
    for i, point in enumerate(health_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(12)

    # Budget audience column
    budget_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(6.75), col_top,
                                        col_width, Inches(4.5))
    budget_bg.fill.solid()
    budget_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    budget_bg.line.color.rgb = SUCCESS

    budget_title = slide.shapes.add_textbox(Inches(7), Inches(1.8), col_width - Inches(0.5), Inches(0.5))
    tf = budget_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Budget Officers & Electeds"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SUCCESS

    budget_content = slide.shapes.add_textbox(Inches(7), Inches(2.4), col_width - Inches(0.5), Inches(3.5))
    tf = budget_content.text_frame
    tf.word_wrap = True

    budget_points = [
        "Lead with financial returns",
        "Show avoided costs clearly",
        "Compare to other investments",
        "\"Every $1 invested returns $X\""
    ]
    for i, point in enumerate(budget_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_PRIMARY
        p.space_after = Pt(12)

    # Bottom note
    note_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.3), Inches(11.833), Inches(0.6))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Same data, different framing. Know your audience."
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER


def add_stay_connected_slide(prs):
    """Add Stay Connected CTA slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.12))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Stay Connected"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.833), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Continue building your health economics skills"
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_SECONDARY

    # Three columns
    col_width = Inches(3.8)
    col_top = Inches(2.3)

    # Column 1
    col1_title = slide.shapes.add_textbox(Inches(0.75), col_top, col_width, Inches(0.5))
    tf = col1_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Practice in Methods Lab"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    col1_content = slide.shapes.add_textbox(Inches(0.75), Inches(2.9), col_width, Inches(1.5))
    tf = col1_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Interactive labs on ROI calculation, building the business case, and presenting to stakeholders."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "caphegroup.org/methods-lab"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.space_before = Pt(12)

    # Column 2
    col2_title = slide.shapes.add_textbox(Inches(4.75), col_top, col_width, Inches(0.5))
    tf = col2_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Get Feedback"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    col2_content = slide.shapes.add_textbox(Inches(4.75), Inches(2.9), col_width, Inches(1.5))
    tf = col2_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Bring your ROI analyses to monthly peer review sessions for expert feedback."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "info@caphegroup.org"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.space_before = Pt(12)

    # Column 3
    col3_title = slide.shapes.add_textbox(Inches(8.75), col_top, col_width, Inches(0.5))
    tf = col3_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Review the Series"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT

    col3_content = slide.shapes.add_textbox(Inches(8.75), Inches(2.9), col_width, Inches(1.5))
    tf = col3_content.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Missed earlier sessions? Recordings available for Feb (Intro) and April (CEA)."
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = "caphegroup.org/programs"
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = PRIMARY
    p2.space_before = Pt(12)

    # Series recap box
    series_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(0.75), Inches(5),
                                         Inches(11.833), Inches(1.8))
    series_box.fill.solid()
    series_box.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF8)
    series_box.line.color.rgb = PRIMARY

    series_title = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.333), Inches(0.4))
    tf = series_title.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 Webinar Series Complete!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PRIMARY

    series_content = slide.shapes.add_textbox(Inches(1), Inches(5.65), Inches(11.333), Inches(1))
    tf = series_content.text_frame
    p = tf.paragraphs[0]
    p.text = "Feb: Intro to Health Economics  •  April: Cost-Effectiveness Basics  •  June: Return on Investment"
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_PRIMARY


def add_closing_slide(prs):
    """Add closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # Key message
    msg_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1))
    tf = msg_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Public health is an investment,"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    msg2_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1))
    tf = msg2_box.text_frame
    p = tf.paragraphs[0]
    p.text = "not just an expense."
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    p.alignment = PP_ALIGN.CENTER

    thanks_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(1.5))
    tf = contact_box.text_frame

    for i, contact in enumerate(["caphegroup.org", "info@caphegroup.org"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = contact
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)


def main():
    """Generate the complete presentation."""
    prs = create_presentation()

    print("Generating slides for: Return on Investment in Public Health")
    print("June 11, 2026 Webinar")
    print("=" * 70)

    # ==========================================================================
    # SLIDE 1: Title
    # ==========================================================================
    add_title_slide(
        prs,
        title="Return on Investment\nin Public Health",
        subtitle="Free Webinar Series",
        date="June 11, 2026 · 12:00 PM PT"
    )
    print("  1. Title slide")

    # ==========================================================================
    # SLIDE 2: The Core Message
    # ==========================================================================
    add_key_concept_slide(
        prs,
        concept="Public Health Is an Investment",
        definition="Every dollar spent on public health generates returns. Healthcare savings. Productivity gains. Reduced emergency costs. Stronger communities.",
        example="The question isn't whether to spend on public health. It's how to show stakeholders the value of what we do.",
        image="roi-concept-investment.png"
    )
    print("  2. Core message")

    # ==========================================================================
    # SLIDE 3: Series Recap
    # ==========================================================================
    add_content_slide(
        prs,
        title="Building on What We've Learned",
        bullets=[
            "February: Economic thinking in public health",
            "  Scarcity, opportunity cost, trade-offs",
            "April: Cost-effectiveness analysis",
            "  ICER, QALYs, decision thresholds",
            "Today: Return on Investment",
            "  Making the financial case for public health"
        ],
        note="Recordings of earlier sessions available at caphegroup.org/programs"
    )
    print("  3. Series recap")

    # ==========================================================================
    # SLIDE 4: Today's Objectives
    # ==========================================================================
    add_content_slide(
        prs,
        title="Today's Objectives",
        bullets=[
            "Understand what ROI means in a public health context",
            "Learn to calculate and interpret ROI",
            "See examples of high-ROI public health investments",
            "Know when ROI is the right metric (and when it isn't)",
            "Build skills to communicate value to budget decision-makers"
        ]
    )
    print("  4. Today's objectives")

    # ==========================================================================
    # SLIDE 5: Section - What is ROI?
    # ==========================================================================
    add_section_slide(prs, "Part 1", "What is Return on Investment?")
    print("  5. Section: What is ROI?")

    # ==========================================================================
    # SLIDE 6: ROI Definition
    # ==========================================================================
    add_key_concept_slide(
        prs,
        concept="Return on Investment (ROI)",
        definition="ROI measures the financial return generated by an investment, expressed as a ratio or percentage. It answers: \"For every dollar we spend, how many dollars do we get back?\"",
        example="An ROI of 5:1 means every $1 invested generates $5 in returns. That's a 400% return on investment.",
        image="roi-concept-multiplier.png"
    )
    print("  6. ROI definition")

    # ==========================================================================
    # SLIDE 7: The ROI Formula
    # ==========================================================================
    add_formula_slide(
        prs,
        title="The ROI Formula",
        formula="ROI = (Returns − Investment) ÷ Investment",
        explanation_lines=[
            "Returns: Total value generated (savings, avoided costs, productivity)",
            "Investment: Total cost of the program or intervention",
            "Result: Expressed as ratio (5:1) or percentage (400%)",
            "Positive ROI means the program pays for itself and more"
        ],
        note="Also expressed as: ROI = Net Benefits ÷ Costs"
    )
    print("  7. The ROI formula")

    # ==========================================================================
    # SLIDE 8: ROI vs CEA
    # ==========================================================================
    add_comparison_table_slide(
        prs,
        title="ROI vs. Cost-Effectiveness: Different Questions",
        headers=["", "ROI", "Cost-Effectiveness"],
        rows=[
            ["Question", "\"Does this pay for itself?\"", "\"What's the best health value?\""],
            ["Units", "Dollars returned per dollar spent", "Cost per health outcome (QALY, life saved)"],
            ["Best for", "Budget justification, elected officials", "Comparing health interventions"],
            ["Limitation", "Requires monetizing all benefits", "Doesn't show if affordable"],
        ],
        note="Both are valuable. Choose based on your audience and question."
    )
    print("  8. ROI vs CEA")

    # ==========================================================================
    # SLIDE 9: What Counts as Returns?
    # ==========================================================================
    add_content_slide(
        prs,
        title="What Counts as \"Returns\"?",
        bullets=[
            "Healthcare cost savings",
            "  Avoided hospitalizations, ER visits, treatments",
            "Productivity gains",
            "  Fewer sick days, longer working lives",
            "Reduced public program costs",
            "  Less Medicaid spending, fewer disability claims",
            "Avoided emergency response costs",
            "  Prevented outbreaks, reduced crisis management",
            "Economic multiplier effects",
            "  Healthier workforce, stronger local economy"
        ],
        image="roi-concept-returns.png"
    )
    print("  9. What counts as returns")

    # ==========================================================================
    # SLIDE 10: Section - Examples
    # ==========================================================================
    add_section_slide(prs, "Part 2", "ROI in Action: California Examples")
    print("  10. Section: Examples")

    # ==========================================================================
    # SLIDE 11: Example 1 - Tobacco Control
    # ==========================================================================
    add_roi_example_slide(
        prs,
        title="California Tobacco Control Program",
        program="Prop 99-funded comprehensive tobacco prevention (1989-2004)",
        investment="$1.8 billion over 15 years",
        returns=[
            "$86 billion in healthcare savings",
            "1+ million life-years saved",
            "33% decline in adult smoking",
            "Model for national and global programs"
        ],
        roi_ratio="$48 returned for every $1 invested",
        source="Lightwood & Glantz, PLOS Medicine 2013",
        image="roi-example-california.png"
    )
    print("  11. Tobacco control example")

    # ==========================================================================
    # SLIDE 12: Example 2 - Childhood Immunization
    # ==========================================================================
    add_roi_example_slide(
        prs,
        title="Childhood Immunization Programs",
        program="Routine childhood vaccination (national estimates, applicable to CA)",
        investment="~$10 billion annually (US)",
        returns=[
            "$69 billion in direct medical savings",
            "$244 billion in total societal savings",
            "21 million hospitalizations prevented",
            "Eliminates or controls 14 diseases"
        ],
        roi_ratio="$7-$24 returned for every $1 invested",
        source="Zhou et al., Pediatrics 2014"
    )
    print("  12. Immunization example")

    # ==========================================================================
    # SLIDE 13: Example 3 - Diabetes Prevention
    # ==========================================================================
    add_roi_example_slide(
        prs,
        title="Diabetes Prevention Programs",
        program="CDC-recognized lifestyle change programs",
        investment="~$500 per participant",
        returns=[
            "58% reduction in diabetes incidence",
            "$2,650 savings per participant (3 years)",
            "Reduced medication costs",
            "Fewer complications, hospitalizations"
        ],
        roi_ratio="$5 returned for every $1 invested",
        source="Diabetes Prevention Program Research Group; CMS Medicare DPP"
    )
    print("  13. Diabetes prevention example")

    # ==========================================================================
    # SLIDE 14: Example 4 - Lead Poisoning Prevention
    # ==========================================================================
    add_roi_example_slide(
        prs,
        title="Lead Poisoning Prevention",
        program="Lead hazard control and screening programs",
        investment="$1-2 billion annually (national)",
        returns=[
            "$50 billion in lifetime earnings preserved",
            "Reduced special education costs",
            "Lower criminal justice involvement",
            "Prevented cognitive impairment"
        ],
        roi_ratio="$17-$221 returned for every $1 invested",
        source="Gould, Environmental Health Perspectives 2009"
    )
    print("  14. Lead prevention example")

    # ==========================================================================
    # SLIDE 15: The Pattern
    # ==========================================================================
    add_content_slide(
        prs,
        title="The Pattern: Prevention Pays",
        bullets=[
            "Upstream investments yield downstream savings",
            "Returns often exceed 5:1 or 10:1",
            "Benefits extend beyond healthcare to productivity and society",
            "Long time horizons favor prevention",
            "  Tobacco: 15-year analysis showed massive returns",
            "  Lead: Lifetime earnings gains from childhood intervention",
            "Public health generates returns that private markets miss"
        ],
        image="roi-concept-prevention-pays.png"
    )
    print("  15. The pattern")

    # ==========================================================================
    # SLIDE 16: Section - Calculating ROI
    # ==========================================================================
    add_section_slide(prs, "Part 3", "Calculating ROI for Your Programs")
    print("  16. Section: Calculating ROI")

    # ==========================================================================
    # SLIDE 17: Steps to Calculate
    # ==========================================================================
    add_content_slide(
        prs,
        title="Steps to Calculate ROI",
        bullets=[
            "1. Define your investment clearly",
            "  All program costs: staff, materials, overhead",
            "2. Identify all relevant returns",
            "  Healthcare savings, productivity, avoided costs",
            "3. Gather data or use credible estimates",
            "  Published studies, local data, reasonable assumptions",
            "4. Choose your time horizon",
            "  Longer horizons capture more returns but add uncertainty",
            "5. Calculate and present with transparency",
            "  Show your assumptions, report ranges"
        ]
    )
    print("  17. Steps to calculate")

    # ==========================================================================
    # SLIDE 18: Data Sources
    # ==========================================================================
    add_content_slide(
        prs,
        title="Where to Find Return Data",
        bullets=[
            "Published ROI studies on similar programs",
            "  County Health Rankings evidence database",
            "  Community Guide (thecommunityguide.org)",
            "Local healthcare cost data",
            "  Hospital discharge data, Medicaid claims",
            "National estimates adapted to local context",
            "  CDC, CMS, peer-reviewed literature",
            "Economic models and calculators",
            "  Prevention ROI calculators (various sources)"
        ]
    )
    print("  18. Data sources")

    # ==========================================================================
    # SLIDE 19: Common Challenges
    # ==========================================================================
    add_content_slide(
        prs,
        title="Common Challenges in ROI Analysis",
        bullets=[
            "Monetizing health outcomes is difficult",
            "  What's a life worth? A prevented disease?",
            "Benefits may accrue to different payers",
            "  You invest, insurance company saves",
            "Time lags between investment and returns",
            "  Prevention today, savings in 10 years",
            "Attribution is imperfect",
            "  How much of the improvement came from our program?",
            "Conservative estimates may understate value",
            "  Societal benefits often excluded"
        ]
    )
    print("  19. Common challenges")

    # ==========================================================================
    # SLIDE 20: When ROI Isn't Enough
    # ==========================================================================
    add_content_slide(
        prs,
        title="When ROI Isn't the Right Metric",
        bullets=[
            "Some programs matter even without positive ROI",
            "  Rare disease services, end-of-life care, equity programs",
            "Not all benefits can be monetized",
            "  Dignity, peace of mind, community trust",
            "ROI favors programs with measurable cost savings",
            "  May undervalue prevention of suffering",
            "Different stakeholders value different things",
            "  Use ROI alongside equity and health impact metrics"
        ],
        note="ROI is a tool, not a decision-maker. It's one input among many."
    )
    print("  20. When ROI isn't enough")

    # ==========================================================================
    # SLIDE 21: Section - Making the Case
    # ==========================================================================
    add_section_slide(prs, "Part 4", "Making the Case to Decision-Makers")
    print("  21. Section: Making the case")

    # ==========================================================================
    # SLIDE 22: Speaking to Your Audience
    # ==========================================================================
    add_audience_framing_slide(prs)
    print("  22. Speaking to your audience")

    # ==========================================================================
    # SLIDE 23: Building the Business Case
    # ==========================================================================
    add_content_slide(
        prs,
        title="Building a Compelling Business Case",
        bullets=[
            "Lead with the bottom line",
            "  \"This program returns $5 for every $1 invested\"",
            "Show where the savings go",
            "  Healthcare system, county budget, employers, families",
            "Use local examples when possible",
            "  \"In our county, this would mean...\"",
            "Acknowledge uncertainty honestly",
            "  \"Estimates range from X to Y\"",
            "Connect to priorities they care about",
            "  Budget stability, workforce health, crisis prevention"
        ]
    )
    print("  23. Building the business case")

    # ==========================================================================
    # SLIDE 24: One-Page Template
    # ==========================================================================
    add_content_slide(
        prs,
        title="The One-Page ROI Summary",
        bullets=[
            "1. The Ask: What you're proposing and the investment needed",
            "2. The Return: Estimated ROI ratio and key benefits",
            "3. The Evidence: Where the numbers come from",
            "4. The Timeline: When returns will be realized",
            "5. The Risk: What happens if we don't invest",
            "6. The Comparison: How this stacks up to alternatives"
        ],
        note="Template available in CAPHE Methods Lab"
    )
    print("  24. One-page template")

    # ==========================================================================
    # SLIDE 25: Stay Connected
    # ==========================================================================
    add_stay_connected_slide(prs)
    print("  25. Stay Connected")

    # ==========================================================================
    # SLIDE 26: Closing
    # ==========================================================================
    add_closing_slide(prs)
    print("  26. Closing slide")

    # Save presentation
    output_dir = Path("/Users/victoriaperez/Projects/CAPHE/07_website/outputs/presentations")
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / "CAPHE_Webinar_ROI_2026-06-11.pptx"

    prs.save(str(output_path))

    print("=" * 70)
    print(f"\n✓ Presentation saved: {output_path}")
    print(f"  Total slides: {len(prs.slides)}")

    return output_path


if __name__ == "__main__":
    main()
